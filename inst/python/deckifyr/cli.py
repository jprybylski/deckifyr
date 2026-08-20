"""The `deckifyr` command-line interface (spec section 11.1).

Every subcommand from the spec is wired up with real argument parsing.
`init`, `validate`, `build`, `preview`, `inspect`, and `schema` do real
work today: `init` copies the bundled minimal example, `validate` loads
and pydantic-validates a project (design + layouts + presentation,
cross-checking layout references and box unit strings), `build` plans
(`deckifyr.plan`) and composes (`deckifyr.pptx`) a `.pptx` and manifest
for `text`/`markdown`/`image`/`shape`/`group`/`table`/`reportifyr`/
`quarto` elements, `preview` builds the same way and then rasterizes
each slide to a PNG via `deckifyr.renderers.preview` (LibreOffice +
PyMuPDF -- requires `soffice` on PATH), `inspect` reports a
presentation.yaml's resolved plan or a built .pptx's real slide/shape
structure (detected by file extension), `schema` dumps a document
type's JSON Schema, and `serve` runs the local web application
(`deckifyr.web.app.create_app`, spec section 12) against `--project`
(default: cwd)/`--presentation` (default: `presentation.yaml`),
blocking until interrupted -- it requires the `web` extra
(`deckifyr[web]`) and raises a plain `DeckifyrError` (exit code 3) with
install guidance if `fastapi`/`uvicorn` aren't importable, the same
"tell the caller what to install, don't crash on ImportError" posture
`deckifyr.renderers.preview`/`deckifyr.renderers.quarto` already take
for their own external dependencies. `skills` (issue #50) exports the
package's own bundled coding-agent skill files (Claude Skills-format
`SKILL.md`s under `deckifyr/skills/`, one for `design.yaml`/
`layouts.yaml` authoring and one for `presentation.yaml` authoring) into
`--directory <target>/<skill-name>/SKILL.md` (default: cwd) -- it never
assumes a `.claude/skills/` layout; the caller picks the target
directory, same `directory` positional + `--force` shape as `init`.

`get`/`set` and the `slide` subcommand group (issue #10) round-trip a
design/layouts/presentation YAML file through `deckifyr.editor`'s pure
dict-manipulation helpers: this module owns reading the file, validating
the edited result against the right `deckifyr.schema` model *before*
writing it back (so a bad edit never corrupts the file on disk), and
turning `deckifyr.editor`'s plain exceptions into `DeckifyrError`s with a
stable code -- `deckifyr.editor`'s own module docstring has the fuller
design writeup. These are real, tested, not stubs.

Exit codes are stable and independent of message wording, per spec
section 11.1:
    0  success
    1  schema/validation failure
    2  argparse usage error (argparse's own default)
    3  I/O failure (missing file, unwritable target, ...)
    4  feature not implemented yet

No subcommand makes network calls (spec section 11.1: "No implicit
network access during a build unless explicitly enabled").
"""

from __future__ import annotations

import argparse
import json
import shutil
import sys
from pathlib import Path
from typing import Any

from pydantic import ValidationError as PydanticValidationError

from deckifyr import editor, projectio
from deckifyr.plan import expand_presentation
from deckifyr.pptx import compose_and_write
from deckifyr.schema.design import DesignDocument
from deckifyr.schema.errors import DeckifyrError, ErrorCode
from deckifyr.schema.layouts import LayoutsDocument
from deckifyr.schema.presentation import PresentationDocument

EXIT_OK = 0
EXIT_VALIDATION_ERROR = 1
EXIT_IO_ERROR = 3
EXIT_NOT_IMPLEMENTED = 4

_EXIT_CODE_BY_ERROR_CODE = {
    ErrorCode.IO: EXIT_IO_ERROR,
    ErrorCode.NOT_IMPLEMENTED: EXIT_NOT_IMPLEMENTED,
}


def _examples_dir() -> Path:
    # inst/python/deckifyr/cli.py -> inst/examples
    return Path(__file__).resolve().parents[2] / "examples" / "minimal-deck"


def _cmd_init(args: argparse.Namespace) -> dict[str, Any]:
    target = Path(args.directory)
    target.mkdir(parents=True, exist_ok=True)
    existing = list(target.iterdir())
    if existing and not args.force:
        raise DeckifyrError(
            f"{target} is not empty (use --force to overwrite)", code=ErrorCode.IO
        )

    created = []
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        source = _examples_dir() / name
        destination = target / name
        shutil.copyfile(source, destination)
        created.append(str(destination))

    return {"directory": str(target), "created": created}


# Bundled Claude Skills-format (SKILL.md) content (issue #50), one
# directory per skill under inst/python/deckifyr/skills/ -- shipped as
# package data (pyproject.toml's package-data), the same "bundled inside
# the package" precedent inst/python/deckifyr/schemas/*.schema.json
# already established for issue #49.
_SKILL_NAMES = ("deckifyr-org-config", "deckifyr-presentation")


def _skills_dir() -> Path:
    return Path(__file__).resolve().parent / "skills"


def _cmd_skills(args: argparse.Namespace) -> dict[str, Any]:
    target = Path(args.directory)

    # Only refuse on the exact destination files, not the whole target
    # directory's non-emptiness (unlike _cmd_init) -- the target may be
    # e.g. `.claude/skills`, which can legitimately already hold other,
    # unrelated skills.
    conflicts = [
        str(target / name / "SKILL.md")
        for name in _SKILL_NAMES
        if (target / name / "SKILL.md").exists()
    ]
    if conflicts and not args.force:
        raise DeckifyrError(
            f"{', '.join(conflicts)} already exist (use --force to overwrite)",
            code=ErrorCode.IO,
        )

    created = []
    for name in _SKILL_NAMES:
        source = _skills_dir() / name / "SKILL.md"
        destination = target / name / "SKILL.md"
        destination.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(source, destination)
        created.append(str(destination))

    return {"directory": str(target), "created": created}


def _cmd_validate(args: argparse.Namespace) -> dict[str, Any]:
    presentation, design, layouts = projectio.load_project(
        Path(args.presentation), strict=args.strict
    )
    return {
        "valid": True,
        "presentation": str(args.presentation),
        "slide_count": len(presentation.slides),
        "layout_count": len(layouts.layouts),
        "schema_version": presentation.deckifyr,
    }


def _cmd_build(args: argparse.Namespace) -> dict[str, Any]:
    presentation_path = Path(args.presentation).resolve()
    # Fail on schema/geometry problems before planning/composing, so
    # `build` on a broken project reports the real validation error
    # rather than a confusing failure downstream.
    presentation, design, layouts = projectio.load_project(
        Path(args.presentation), strict=args.strict
    )

    project_root = presentation_path.parent
    resolved_slides = expand_presentation(
        presentation, design, layouts, strict=args.strict
    )
    # `keep_preview_pdf=True`: harmless when `build.previews` is off (no
    # preview render happens at all, so `result.preview_pdf_path` stays
    # `None`), and when a preview render does happen this build is already
    # paying LibreOffice's PDF-conversion cost to make the PNGs -- keeping
    # that intermediate PDF alongside the built `.pptx` (issue #32) is
    # free, the same reasoning `_cmd_preview` below already uses for its
    # own `keep_preview_pdf=True`.
    result = compose_and_write(
        presentation,
        design,
        resolved_slides,
        project_root=project_root,
        presentation_path=presentation_path,
        design_path=(project_root / presentation.design.base).resolve(),
        layouts_path=(project_root / presentation.layouts).resolve(),
        keep_preview_pdf=True,
    )

    return {
        "output": str(result.output_path),
        "manifest": str(result.manifest_path) if result.manifest_path else None,
        "slide_count": result.slide_count,
        "warning_count": len(result.warnings),
        "previews": [str(p) for p in result.preview_paths],
        "preview_pdf": str(result.preview_pdf_path) if result.preview_pdf_path else None,
    }


def _cmd_get(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.file)
    data = projectio.read_yaml(path)
    try:
        value = editor.get_value(data, args.path)
    except editor.PathError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.PATH_NOT_FOUND) from exc
    return {"file": str(path), "path": args.path, "value": value}


def _parse_set_value(raw: str) -> Any:
    # JSON, not YAML: a YAML scalar parser would treat a bare hex color
    # ("#123456", extremely common in design.yaml) as a comment opener
    # and silently parse it as `None` -- confirmed the hard way while
    # smoke-testing this command by hand. JSON has no comment syntax at
    # all, so `json.loads` either parses `raw` unambiguously as a
    # number/bool/null/array/object (matching --elements-json's own
    # vocabulary) or raises -- at which point `raw` was never valid JSON
    # to begin with, so it's used as a literal string. This means an
    # ordinary bare word/hex color/font name needs no quoting on the
    # command line, while `'"12pt"'`/`true`/`[1, 2]` still work when a
    # caller actually wants a typed value.
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        return raw


def _cmd_set(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.file)
    data = projectio.read_yaml(path)

    value: Any = args.value if args.string else _parse_set_value(args.value)

    try:
        editor.set_value(data, args.path, value)
    except editor.PathError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.PATH_NOT_FOUND) from exc

    doc_type = projectio.resolve_document_type(args.type, data)
    if doc_type == "presentation":
        extra = projectio.validate_and_write_presentation(path, data)
    else:
        model = projectio.DOCUMENT_MODELS[doc_type]
        try:
            model.model_validate(data)
        except PydanticValidationError as exc:
            raise DeckifyrError(
                projectio.format_pydantic_error(exc, str(path)),
                code=ErrorCode.SCHEMA_VALIDATION,
            ) from exc
        projectio.write_yaml(path, data)
        extra = {}

    return {"file": str(path), "path": args.path, "type": doc_type, **extra}


def _cmd_slide_list(args: argparse.Namespace) -> dict[str, Any]:
    data = projectio.load_presentation_raw(Path(args.presentation))
    return {"presentation": args.presentation, "slides": editor.list_slides(data)}


def _cmd_slide_add(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.presentation)
    data = projectio.load_presentation_raw(path)
    elements = (
        projectio.parse_json_arg(args.elements_json, "--elements-json")
        if args.elements_json is not None
        else None
    )
    try:
        editor.add_slide(
            data,
            id=args.id,
            layout=args.layout,
            elements=elements,
            notes=args.notes,
            index=args.index,
            after=args.after,
            before=args.before,
        )
    except editor.DuplicateSlideIdError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.CONTENT_VALIDATION) from exc
    except editor.SlideNotFoundError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.REFERENCE_NOT_FOUND) from exc
    except editor.AmbiguousPlacementError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.CONTENT_VALIDATION) from exc
    return projectio.validate_and_write_presentation(path, data)


def _cmd_slide_remove(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.presentation)
    data = projectio.load_presentation_raw(path)
    try:
        editor.remove_slide(data, args.id)
    except editor.SlideNotFoundError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.REFERENCE_NOT_FOUND) from exc
    return projectio.validate_and_write_presentation(path, data)


def _cmd_slide_update(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.presentation)
    data = projectio.load_presentation_raw(path)

    kwargs: dict[str, Any] = {}
    if args.no_layout:
        kwargs["layout"] = None
    elif args.layout is not None:
        kwargs["layout"] = args.layout
    if args.clear_notes:
        kwargs["notes"] = None
    elif args.notes is not None:
        kwargs["notes"] = args.notes
    if args.elements_json is not None:
        kwargs["elements"] = projectio.parse_json_arg(args.elements_json, "--elements-json")

    try:
        editor.update_slide(data, args.id, **kwargs)
    except editor.SlideNotFoundError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.REFERENCE_NOT_FOUND) from exc
    return projectio.validate_and_write_presentation(path, data)


def _cmd_slide_move(args: argparse.Namespace) -> dict[str, Any]:
    path = Path(args.presentation)
    data = projectio.load_presentation_raw(path)
    try:
        editor.move_slide(
            data, args.id, index=args.index, after=args.after, before=args.before
        )
    except editor.SlideNotFoundError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.REFERENCE_NOT_FOUND) from exc
    except editor.AmbiguousPlacementError as exc:
        raise DeckifyrError(str(exc), code=ErrorCode.CONTENT_VALIDATION) from exc
    return projectio.validate_and_write_presentation(path, data)


def _parse_slides_arg(raw: str | None) -> list[int] | None:
    """Parses `--slides`' comma-separated 1-indexed list (`"1,3"`) into
    `[1, 3]`, or `None` for "not given" (render every slide, unchanged
    from before this flag existed).
    """
    if raw is None:
        return None
    try:
        return [int(token.strip()) for token in raw.split(",") if token.strip()]
    except ValueError as exc:
        raise DeckifyrError(
            f"--slides must be a comma-separated list of integers, got {raw!r}",
            code=ErrorCode.CONTENT_VALIDATION,
        ) from exc


def _cmd_preview(args: argparse.Namespace) -> dict[str, Any]:
    presentation_path = Path(args.presentation).resolve()
    presentation, design, layouts = projectio.load_project(
        Path(args.presentation), strict=args.strict
    )

    project_root = presentation_path.parent
    resolved_slides = expand_presentation(
        presentation, design, layouts, strict=args.strict
    )
    # `force_previews=True`: an explicit `deckifyr preview` invocation
    # always renders, regardless of the project's own `build.previews`
    # flag -- see `compose_and_write`'s own docstring note on this.
    # `keep_preview_pdf=True`: this command already pays the LibreOffice
    # PDF-conversion cost, so keeping the intermediate PDF around (issue
    # #27's embedded-PDF-viewer support) is free -- `_cmd_build` above
    # does the same (issue #32) whenever `build.previews` actually
    # triggers a render.
    result = compose_and_write(
        presentation,
        design,
        resolved_slides,
        project_root=project_root,
        presentation_path=presentation_path,
        design_path=(project_root / presentation.design.base).resolve(),
        layouts_path=(project_root / presentation.layouts).resolve(),
        force_previews=True,
        preview_slides=_parse_slides_arg(args.slides),
        keep_preview_pdf=True,
    )

    return {
        "output": str(result.output_path),
        "previews": [str(p) for p in result.preview_paths],
        "preview_pdf": str(result.preview_pdf_path) if result.preview_pdf_path else None,
        "slide_count": result.slide_count,
    }


def _inspect_presentation(path: Path, *, strict: bool) -> dict[str, Any]:
    presentation, design, layouts = projectio.load_project(path, strict=strict)
    resolved_slides = expand_presentation(presentation, design, layouts, strict=strict)
    return {
        "target": "presentation",
        "path": str(path),
        "schema_version": presentation.deckifyr,
        "slide_count": len(resolved_slides),
        "layout_count": len(layouts.layouts),
        "status_indicator": presentation.status_indicator,
        "slides": [
            {
                "id": slide.id,
                "element_count": len(slide.elements),
                "element_types": sorted({element.type for element in slide.elements}),
                "has_notes": slide.notes is not None,
            }
            for slide in resolved_slides
        ],
    }


def _inspect_pptx(path: Path) -> dict[str, Any]:
    # Imported lazily so `deckifyr inspect some.yaml` never pays for
    # importing python-pptx's own presentation-reading machinery, the
    # same lazy-import posture `deckifyr.resolvers.table` already takes
    # for pyarrow.
    from pptx import Presentation as PptxPresentation

    try:
        prs = PptxPresentation(str(path))
    except Exception as exc:  # python-pptx raises a mix of exception types
        raise DeckifyrError(
            f"{path}: not a readable .pptx package: {exc}", code=ErrorCode.IO
        ) from exc

    slides = []
    for index, slide in enumerate(prs.slides):
        shapes = [
            {
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "has_text_frame": shape.has_text_frame,
                "rotation": shape.rotation,
            }
            for shape in slide.shapes
        ]
        has_notes = (
            slide.has_notes_slide
            and slide.notes_slide.notes_text_frame.text.strip() != ""
        )
        slides.append(
            {
                "index": index,
                "shape_count": len(shapes),
                "shapes": shapes,
                "has_notes": has_notes,
            }
        )

    manifest_summary = None
    manifest_path = path.with_name(path.stem + ".manifest.json")
    if manifest_path.is_file():
        try:
            manifest_data = json.loads(manifest_path.read_text())
        except (OSError, json.JSONDecodeError):
            manifest_data = None
        if isinstance(manifest_data, dict):
            manifest_summary = {
                "path": str(manifest_path),
                "deckifyr_version": manifest_data.get("deckifyr_version"),
                "slide_count": manifest_data.get("slide_count"),
                "warnings": manifest_data.get("warnings", []),
            }

    return {
        "target": "pptx",
        "path": str(path),
        "slide_width_emu": prs.slide_width,
        "slide_height_emu": prs.slide_height,
        "slide_count": len(slides),
        "slides": slides,
        "manifest": manifest_summary,
    }


def _cmd_inspect(args: argparse.Namespace) -> dict[str, Any]:
    target = Path(args.target)
    if not target.is_file():
        raise DeckifyrError(f"target not found: {target}", code=ErrorCode.IO)

    suffix = target.suffix.lower()
    if suffix in (".yaml", ".yml"):
        return _inspect_presentation(target, strict=args.strict)
    if suffix == ".pptx":
        return _inspect_pptx(target)
    raise DeckifyrError(
        f"cannot infer target type from extension {suffix!r} -- expected "
        "a presentation .yaml/.yml or a built .pptx",
        code=ErrorCode.IO,
    )


def _cmd_serve(args: argparse.Namespace) -> dict[str, Any]:
    # Imported lazily: the `web` extra (fastapi/uvicorn) is optional, and
    # every other subcommand must keep working without it installed --
    # the same posture `deckifyr.renderers.preview`/`deckifyr.renderers
    # .quarto` already take for their own external dependencies.
    try:
        import uvicorn
    except ImportError as exc:
        raise DeckifyrError(
            "the 'web' extra is not installed -- run `pip install "
            "deckifyr[web]` (or `uv sync --extra web`)",
            code=ErrorCode.IO,
        ) from exc

    from deckifyr.web.app import create_app

    project_root = Path(args.project or ".").resolve()
    app = create_app(project_root, args.presentation, launcher=args.launcher)
    # Blocks until interrupted (Ctrl-C / SIGINT/SIGTERM) -- uvicorn's own
    # default graceful shutdown handling, nothing custom here.
    uvicorn.run(app, host=args.host, port=args.port)
    return {"status": "stopped"}


_SCHEMA_MODELS = {
    "design": DesignDocument,
    "layouts": LayoutsDocument,
    "presentation": PresentationDocument,
}


def _cmd_schema(args: argparse.Namespace) -> dict[str, Any]:
    model = _SCHEMA_MODELS[args.document]
    return model.model_json_schema()


def _build_arg_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="deckifyr")
    parser.add_argument(
        "--json", action="store_true", help="emit structured JSON output"
    )
    subparsers = parser.add_subparsers(dest="command", required=True)

    init_parser = subparsers.add_parser(
        "init", help="scaffold a new project from the bundled minimal example"
    )
    init_parser.add_argument("directory", nargs="?", default=".")
    init_parser.add_argument("--force", action="store_true")
    init_parser.set_defaults(handler=_cmd_init)

    skills_parser = subparsers.add_parser(
        "skills",
        help="export bundled coding-agent skill files (SKILL.md) for authoring "
        "design/layouts/presentation YAML",
    )
    skills_parser.add_argument("directory", nargs="?", default=".")
    skills_parser.add_argument("--force", action="store_true")
    skills_parser.set_defaults(handler=_cmd_skills)

    def add_strict_flag(p: argparse.ArgumentParser) -> None:
        group = p.add_mutually_exclusive_group()
        group.add_argument(
            "--strict", dest="strict", action="store_true", default=True
        )
        group.add_argument("--warn-only", dest="strict", action="store_false")

    validate_parser = subparsers.add_parser(
        "validate", help="validate a presentation.yaml and its referenced files"
    )
    validate_parser.add_argument("presentation")
    add_strict_flag(validate_parser)
    validate_parser.set_defaults(handler=_cmd_validate)

    build_parser = subparsers.add_parser("build", help="build a .pptx and manifest")
    build_parser.add_argument("presentation")
    add_strict_flag(build_parser)
    build_parser.set_defaults(handler=_cmd_build)

    preview_parser = subparsers.add_parser(
        "preview", help="render each slide to a PNG (requires LibreOffice on PATH)"
    )
    preview_parser.add_argument("presentation")
    preview_parser.add_argument(
        "--slides",
        default=None,
        help="comma-separated 1-indexed slide numbers to render (default: all), e.g. 1,3",
    )
    add_strict_flag(preview_parser)
    preview_parser.set_defaults(handler=_cmd_preview)

    inspect_parser = subparsers.add_parser(
        "inspect", help="inspect a presentation.yaml's resolved plan or a built .pptx"
    )
    inspect_parser.add_argument("target")
    add_strict_flag(inspect_parser)
    inspect_parser.set_defaults(handler=_cmd_inspect)

    schema_parser = subparsers.add_parser("schema", help="print a document type's JSON Schema")
    schema_parser.add_argument("document", choices=sorted(_SCHEMA_MODELS))
    schema_parser.set_defaults(handler=_cmd_schema)

    get_parser = subparsers.add_parser(
        "get", help="read a config value from a design/layouts/presentation YAML file"
    )
    get_parser.add_argument("file")
    get_parser.add_argument(
        "path", help="dotted path, e.g. colors.primary or slides[0].notes"
    )
    get_parser.set_defaults(handler=_cmd_get)

    set_parser = subparsers.add_parser(
        "set", help="write a config value into a design/layouts/presentation YAML file"
    )
    set_parser.add_argument("file")
    set_parser.add_argument("path")
    set_parser.add_argument("value")
    set_parser.add_argument(
        "--string",
        action="store_true",
        help=(
            "treat value as a literal string, never as JSON -- needed only for a "
            "value that would otherwise parse as a number/bool/null/array/object "
            "(e.g. the literal text 'true' or 'null')"
        ),
    )
    set_parser.add_argument(
        "--type",
        choices=["auto", "design", "layouts", "presentation"],
        default="auto",
        help="which schema to validate the edited document against (default: auto-detect)",
    )
    set_parser.set_defaults(handler=_cmd_set)

    def add_placement_flags(p: argparse.ArgumentParser) -> None:
        group = p.add_mutually_exclusive_group()
        group.add_argument("--index", type=int, default=None, help="0-based position")
        group.add_argument("--after", default=None, help="place immediately after this slide id")
        group.add_argument("--before", default=None, help="place immediately before this slide id")

    slide_parser = subparsers.add_parser(
        "slide", help="add, remove, update, move, or list presentation.yaml's slides"
    )
    slide_subparsers = slide_parser.add_subparsers(dest="slide_command", required=True)

    slide_list_parser = slide_subparsers.add_parser("list", help="list slides in order")
    slide_list_parser.add_argument("presentation")
    slide_list_parser.set_defaults(handler=_cmd_slide_list)

    slide_add_parser = slide_subparsers.add_parser("add", help="add a new slide")
    slide_add_parser.add_argument("presentation")
    slide_add_parser.add_argument("--id", required=True, help="the new slide's unique id")
    slide_add_parser.add_argument(
        "--layout", default=None, help="layout name (omit for a freeform 'layout: null' slide)"
    )
    slide_add_parser.add_argument("--notes", default=None, help="speaker notes text")
    slide_add_parser.add_argument(
        "--elements-json",
        default=None,
        help="JSON object/array to use as the slide's 'elements' block",
    )
    add_placement_flags(slide_add_parser)
    slide_add_parser.set_defaults(handler=_cmd_slide_add)

    slide_remove_parser = slide_subparsers.add_parser("remove", help="remove a slide by id")
    slide_remove_parser.add_argument("presentation")
    slide_remove_parser.add_argument("id")
    slide_remove_parser.set_defaults(handler=_cmd_slide_remove)

    slide_update_parser = slide_subparsers.add_parser(
        "update", help="update an existing slide's layout, notes, or elements"
    )
    slide_update_parser.add_argument("presentation")
    slide_update_parser.add_argument("id")
    layout_group = slide_update_parser.add_mutually_exclusive_group()
    layout_group.add_argument("--layout", default=None, help="new layout name")
    layout_group.add_argument(
        "--no-layout", action="store_true", help="clear the layout (freeform 'layout: null')"
    )
    notes_group = slide_update_parser.add_mutually_exclusive_group()
    notes_group.add_argument("--notes", default=None, help="new speaker notes text")
    notes_group.add_argument("--clear-notes", action="store_true", help="remove speaker notes")
    slide_update_parser.add_argument(
        "--elements-json",
        default=None,
        help="JSON object/array to replace the slide's 'elements' block",
    )
    slide_update_parser.set_defaults(handler=_cmd_slide_update)

    slide_move_parser = slide_subparsers.add_parser("move", help="reorder a slide")
    slide_move_parser.add_argument("presentation")
    slide_move_parser.add_argument("id")
    add_placement_flags(slide_move_parser)
    slide_move_parser.set_defaults(handler=_cmd_slide_move)

    serve_parser = subparsers.add_parser(
        "serve", help="run the local web application (requires the 'web' extra)"
    )
    serve_parser.add_argument("--host", default="127.0.0.1")
    serve_parser.add_argument("--port", type=int, default=8000)
    serve_parser.add_argument(
        "--project", default=None, help="project directory (default: current directory)"
    )
    serve_parser.add_argument(
        "--presentation",
        default="presentation.yaml",
        help="presentation.yaml path, relative to --project (default: presentation.yaml)",
    )
    serve_parser.add_argument(
        "--launcher",
        choices=["cli", "r"],
        default="cli",
        help=(
            "who's launching this (default: cli) -- surfaced via GET /api/health "
            "so the frontend's 'no project found' screen can show deckifyr-CLI vs "
            "R-facade instructions; R/serve.R's deck_serve() passes --launcher r"
        ),
    )
    serve_parser.set_defaults(handler=_cmd_serve)

    return parser


def main(argv: list[str] | None = None) -> int:
    parser = _build_arg_parser()
    args = parser.parse_args(argv)

    try:
        result = args.handler(args)
    except DeckifyrError as exc:
        exit_code = _EXIT_CODE_BY_ERROR_CODE.get(exc.code, EXIT_VALIDATION_ERROR)
        if args.json:
            # Deliberately stderr, not stdout, on the error path: the R
            # facade invokes this CLI through pyro::run_python_script(),
            # which uses processx::run(error_on_status = TRUE) and
            # discards the captured stdout/stderr return value whenever
            # the process exits non-zero, replacing it with a generic
            # "<script> failed." error. R/run-python.R works around this
            # by passing its own stderr_callback to capture output as it
            # streams (before the exit-status check fires), which only
            # works if the diagnostic actually went to stderr. Keeping
            # errors on stderr and success on stdout is also just the
            # more conventional CLI split. See R/run-python.R for the R
            # side of this handshake -- the two must not drift apart.
            print(json.dumps({"status": "error", **exc.to_dict()}, indent=2), file=sys.stderr)
        else:
            print(f"ERROR [{exc.code}]: {exc}", file=sys.stderr)
        return exit_code

    if args.json:
        print(json.dumps({"status": "ok", **result}, indent=2))
    elif args.command == "schema":
        # A JSON Schema has no nicer "human-readable" form than its own
        # JSON, so print it as-is regardless of --json.
        print(json.dumps(result, indent=2))
    else:
        for key, value in result.items():
            print(f"{key}: {value}")
    return EXIT_OK


if __name__ == "__main__":
    raise SystemExit(main())
