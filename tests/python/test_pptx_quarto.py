"""Tests for `deckifyr.pptx.compose`'s `quarto` element support (spec
section 8.1, issue #3).

Fast tests monkeypatch `deckifyr.pptx.compose.QuartoResolver` with a
fake that returns a canned `QuartoArtifact` -- they exercise the
compositor's own logic (native-text placement via `_add_text_shape`,
alt-text enforcement for image modes, manifest `render_mode`/
`editability` recording the *resolved* mode, group recursion) without
needing a real `quarto` binary, so they always run. `test_quarto_*_end_to_end`
at the bottom instead runs the real `quarto`/Typst/PyMuPDF pipeline and
skips cleanly when `quarto` isn't on PATH, mirroring
`test_renderers_quarto.py`'s own skip pattern.
"""

from __future__ import annotations

import importlib
import json
import shutil
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation

from deckifyr.plan import expand_presentation
from deckifyr.pptx.compose import compose_and_write
from deckifyr.resolvers import QuartoArtifact, ResolvedContent
from deckifyr.schema.design import DesignDocument, Fonts, SlideSize, TextStyle
from deckifyr.schema.errors import ContentValidationError
from deckifyr.schema.layouts import Box, Element, Layout, LayoutsDocument
from deckifyr.schema.presentation import (
    BuildConfig,
    DesignRef,
    Metadata,
    PresentationDocument,
    Slide,
)

# `deckifyr.pptx`'s own `__init__.py` does `from deckifyr.pptx.compose
# import ..., compose, ...`, which shadows the `deckifyr.pptx.compose`
# *submodule* attribute on the `deckifyr.pptx` package with the
# `compose()` function of the same name -- `import deckifyr.pptx.compose
# as X` resolves through that (now-shadowed) attribute, so it binds `X`
# to the function, not the module. `importlib.import_module` reads
# `sys.modules` directly instead, sidestepping the shadowing.
compose_module = importlib.import_module("deckifyr.pptx.compose")

requires_quarto = pytest.mark.skipif(
    shutil.which("quarto") is None, reason="quarto binary not found on PATH"
)


def _design() -> DesignDocument:
    return DesignDocument(
        deckifyr="0.1",
        slide=SlideSize(width="10in", height="7.5in"),
        fonts=Fonts(body="Arial", heading="Arial"),
        text_styles={"title": TextStyle(font="heading", size="24pt", bold=True, color="#000000")},
    )


def _presentation(*, elements: list[Element]) -> PresentationDocument:
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[Slide(id="s1", layout=None, elements=elements)],
    )


@pytest.fixture
def project(tmp_path):
    (tmp_path / "design.yaml").write_text("design")
    (tmp_path / "layouts.yaml").write_text("layouts")
    (tmp_path / "presentation.yaml").write_text("presentation")
    (tmp_path / "frag.qmd").write_text("placeholder -- resolution is faked in these tests\n")
    return tmp_path


def _build(project, presentation, design=None):
    layouts = LayoutsDocument(deckifyr="0.1", layouts={"blank": Layout()})
    resolved = expand_presentation(presentation, design or _design(), layouts, strict=True)
    return compose_and_write(
        presentation,
        design or _design(),
        resolved,
        project_root=project,
        presentation_path=project / "presentation.yaml",
        design_path=project / "design.yaml",
        layouts_path=project / "layouts.yaml",
    )


class _FakeQuartoResolver:
    """Stands in for `deckifyr.resolvers.QuartoResolver` -- returns a
    fixed artifact instead of shelling out to a real `quarto` binary.
    """

    artifact: QuartoArtifact | None = None

    def __init__(self, *, config=None) -> None:
        pass

    def resolve(self, value, context, *, requested_render_mode="auto", font=None, text_color=None):
        return ResolvedContent(value=type(self).artifact)


@pytest.fixture
def fake_quarto(monkeypatch):
    monkeypatch.setattr(compose_module, "QuartoResolver", _FakeQuartoResolver)
    return _FakeQuartoResolver


def _quarto_element(**overrides) -> Element:
    data = dict(
        id="frag",
        type="quarto",
        source="frag.qmd",
        box=Box(x="1in", y="1in", width="4in", height="2in"),
    )
    data.update(overrides)
    return Element(**data)


def test_native_mode_places_text_via_markdown_pipeline(project, fake_quarto):
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="native", markdown="**bold** and plain"
    )
    result = _build(project, _presentation(elements=[_quarto_element()]))

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (shape,) = list(slide.shapes)
    assert shape.name == "frag"
    runs = shape.text_frame.paragraphs[0].runs
    assert runs[0].text == "bold" and runs[0].font.bold is True
    assert runs[1].text == " and plain" and runs[1].font.bold is False

    manifest = json.loads(result.manifest_path.read_text())
    (entry,) = manifest["elements"]
    assert entry["render_mode"] == "native"
    assert entry["editability"] == "fully_editable"
    assert entry["resolved_path"].endswith("frag.qmd")


def test_svg_mode_is_rejected_for_pptx_composition(project, fake_quarto):
    # python-pptx cannot embed SVG at all (confirmed against a real
    # render -- see deckifyr.renderers.quarto's module docstring), so an
    # element that resolves to `svg` must fail clearly rather than crash
    # deep inside `_place_picture`'s Pillow-based sizing.
    image_path = project / "rendered.svg"
    image_path.write_text("<svg></svg>")
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="svg", image_path=image_path, image_format="svg"
    )
    with pytest.raises(ContentValidationError, match="render_mode: svg"):
        _build(project, _presentation(elements=[_quarto_element(alt_text="x", render_mode="svg")]))
    assert not image_path.exists()


def test_png_mode_requires_alt_text(project, fake_quarto):
    image_path = project / "rendered.png"
    Image.new("RGB", (10, 10), color="green").save(image_path)
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="png", image_path=image_path, image_format="png"
    )
    with pytest.raises(ContentValidationError, match="alt_text"):
        _build(project, _presentation(elements=[_quarto_element(alt_text=None)]))


def test_png_mode_places_picture_and_records_editability(project, fake_quarto):
    image_path = project / "rendered.png"
    Image.new("RGB", (200, 100), color="blue").save(image_path)
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="png", image_path=image_path, image_format="png"
    )
    result = _build(
        project,
        _presentation(elements=[_quarto_element(alt_text="a rendered fragment", fit="contain")]),
    )

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (picture,) = list(slide.shapes)
    assert picture.name == "frag"

    manifest = json.loads(result.manifest_path.read_text())
    (entry,) = manifest["elements"]
    assert entry["render_mode"] == "png"
    assert entry["editability"] == "rendered_graphic"


def test_manifest_records_resolved_mode_not_the_literal_auto(project, fake_quarto):
    image_path = project / "rendered.png"
    Image.new("RGB", (10, 10), color="green").save(image_path)
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="png", image_path=image_path, image_format="png"
    )
    result = _build(
        project,
        _presentation(elements=[_quarto_element(alt_text="a fragment", render_mode="auto")]),
    )
    manifest = json.loads(result.manifest_path.read_text())
    (entry,) = manifest["elements"]
    assert entry["render_mode"] == "png"


def test_quarto_element_inside_a_group_is_composed(project, fake_quarto):
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="native", markdown="grouped text"
    )
    group = Element(
        id="card",
        type="group",
        box=Box(x="0in", y="0in", width="4in", height="3in"),
        elements=[_quarto_element(id="frag")],
    )
    result = _build(project, _presentation(elements=[group]))

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (group_shape,) = list(slide.shapes)
    assert [s.name for s in group_shape.shapes] == ["frag"]


def test_quarto_element_requires_pathlib_image_cleanup(project, fake_quarto, tmp_path):
    image_path = project / "rendered.png"
    Image.new("RGB", (10, 10), color="red").save(image_path)
    fake_quarto.artifact = QuartoArtifact(
        path=project / "frag.qmd", render_mode="png", image_path=image_path, image_format="png"
    )
    _build(project, _presentation(elements=[_quarto_element(alt_text="x")]))
    assert not image_path.exists()


# ---------------------------------------------------------------------------
# Real end-to-end (requires a real quarto binary)
# ---------------------------------------------------------------------------


@requires_quarto
def test_quarto_native_end_to_end(tmp_path):
    (tmp_path / "design.yaml").write_text("design")
    (tmp_path / "layouts.yaml").write_text("layouts")
    (tmp_path / "presentation.yaml").write_text("presentation")
    (tmp_path / "frag.qmd").write_text("The answer is **42**.\n")

    presentation = _presentation(
        elements=[_quarto_element(render_mode="native")]
    )
    result = _build(tmp_path, presentation)

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (shape,) = list(slide.shapes)
    text = "".join(run.text for para in shape.text_frame.paragraphs for run in para.runs)
    assert "42" in text


@requires_quarto
def test_quarto_auto_math_end_to_end(tmp_path):
    pytest.importorskip("pymupdf")
    (tmp_path / "design.yaml").write_text("design")
    (tmp_path / "layouts.yaml").write_text("layouts")
    (tmp_path / "presentation.yaml").write_text("presentation")
    (tmp_path / "frag.qmd").write_text("$$x^2 + y^2 = z^2$$\n")

    presentation = _presentation(
        elements=[_quarto_element(alt_text="the pythagorean identity")]
    )
    result = _build(tmp_path, presentation)

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (picture,) = list(slide.shapes)
    assert picture.shape_type is not None

    manifest = json.loads(result.manifest_path.read_text())
    (entry,) = manifest["elements"]
    assert entry["render_mode"] == "png"
    assert entry["editability"] == "rendered_graphic"
