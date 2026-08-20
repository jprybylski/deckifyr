import json
import shutil
import subprocess
from pathlib import Path

import pytest
import yaml
from pptx import Presentation

from deckifyr.cli import EXIT_OK, EXIT_VALIDATION_ERROR, _SKILL_NAMES, main

requires_soffice = pytest.mark.skipif(
    shutil.which("soffice") is None, reason="soffice binary not found on PATH"
)
requires_git = pytest.mark.skipif(
    shutil.which("git") is None, reason="git binary not found on PATH"
)


def test_validate_exits_ok_on_minimal_deck(minimal_deck_dir, capsys):
    exit_code = main(["--json", "validate", str(minimal_deck_dir / "presentation.yaml")])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["status"] == "ok"
    assert output["valid"] is True
    assert output["slide_count"] == 2


def test_build_writes_pptx_and_manifest(minimal_deck_dir, tmp_path, capsys):
    # Copy the fixture into tmp_path rather than building in place, so
    # this test doesn't leave a build/ directory inside the repo's
    # bundled example (the same fixture `deckifyr init` scaffolds from).
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)

    exit_code = main(["--json", "build", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["status"] == "ok"
    assert output["slide_count"] == 2

    prs = Presentation(output["output"])
    shape_names = {shape.name for slide in prs.slides for shape in slide.shapes}
    assert shape_names == {"deck-title", "title", "content"}

    manifest = json.loads(Path(output["manifest"]).read_text())
    assert manifest["slide_count"] == 2
    assert {"deckifyr_version", "elements", "input_files", "output"} <= manifest.keys()


@requires_soffice
def test_build_with_previews_enabled_also_keeps_the_pdf(minimal_deck_dir, tmp_path, capsys):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)
    presentation_path = tmp_path / "presentation.yaml"
    text = presentation_path.read_text()
    assert "previews: false" in text
    presentation_path.write_text(text.replace("previews: false", "previews: true"))

    exit_code = main(["--json", "build", str(presentation_path)])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert len(output["previews"]) == 2
    for preview_path in output["previews"]:
        assert Path(preview_path).is_file()
    # Issue #32: an ordinary build with `build.previews: true` now also
    # keeps the intermediate PDF, not just `deckifyr preview`.
    assert output["preview_pdf"] is not None
    assert Path(output["preview_pdf"]).is_file()


def test_validate_reports_missing_file(capsys):
    exit_code = main(["--json", "validate", "does/not/exist.yaml"])
    assert exit_code != EXIT_OK
    output = json.loads(capsys.readouterr().err)
    assert output["code"] == "E_IO"


def test_validate_rejects_unknown_layout_reference(tmp_path, capsys):
    (tmp_path / "design.yaml").write_text(
        "deckifyr: '0.1'\n"
        "slide: {width: 1in, height: 1in}\n"
        "fonts: {body: Arial, heading: Arial}\n"
    )
    (tmp_path / "layouts.yaml").write_text("deckifyr: '0.1'\nlayouts: {blank: {}}\n")
    (tmp_path / "presentation.yaml").write_text(
        "deckifyr: '0.1'\n"
        "design: {base: design.yaml}\n"
        "layouts: layouts.yaml\n"
        "metadata: {title: Test}\n"
        "build: {output: build/out.pptx}\n"
        "slides:\n"
        "  - id: only-slide\n"
        "    layout: does-not-exist\n"
        "    elements: {}\n"
    )
    exit_code = main(["--json", "validate", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_VALIDATION_ERROR
    output = json.loads(capsys.readouterr().err)
    assert output["code"] == "E_REFERENCE_NOT_FOUND"


def test_build_composes_a_color_derived_from_another_token(tmp_path, capsys):
    import colorsys

    (tmp_path / "design.yaml").write_text(
        "deckifyr: '0.1'\n"
        "slide: {width: 1in, height: 1in}\n"
        "fonts: {body: Arial, heading: Arial}\n"
        "colors:\n"
        "  primary: '#2457A6'\n"
        "  secondary:\n"
        "    base: primary\n"
        "    darken: 0.2\n"
        "text_styles:\n"
        "  derived:\n"
        "    font: body\n"
        "    size: 12pt\n"
        "    color: secondary\n"
    )
    (tmp_path / "layouts.yaml").write_text("deckifyr: '0.1'\nlayouts: {blank: {}}\n")
    (tmp_path / "presentation.yaml").write_text(
        "deckifyr: '0.1'\n"
        "design: {base: design.yaml}\n"
        "layouts: layouts.yaml\n"
        "metadata: {title: Test}\n"
        "build: {output: build/out.pptx}\n"
        "slides:\n"
        "  - id: only-slide\n"
        "    layout: null\n"
        "    elements:\n"
        "      - id: label\n"
        "        type: text\n"
        "        value: hi\n"
        "        style: derived\n"
        "        box: {x: 0in, y: 0in, width: 1in, height: 1in}\n"
    )

    exit_code = main(["--json", "build", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)

    r, g, b = (int("2457A6"[i : i + 2], 16) / 255.0 for i in (0, 2, 4))
    h, l, s = colorsys.rgb_to_hls(r, g, b)
    expected_rgb = colorsys.hls_to_rgb(h, max(0.0, l - 0.2), s)
    expected_hex = "".join(f"{round(c * 255):02X}" for c in expected_rgb)

    prs = Presentation(output["output"])
    (slide,) = list(prs.slides)
    (shape,) = list(slide.shapes)
    run = shape.text_frame.paragraphs[0].runs[0]
    assert str(run.font.color.rgb) == expected_hex


def test_schema_command_prints_json_schema(capsys):
    exit_code = main(["schema", "presentation"])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["title"] == "PresentationDocument"


def test_init_scaffolds_a_new_project(tmp_path):
    target = tmp_path / "new-project"
    exit_code = main(["--json", "init", str(target)])
    assert exit_code == EXIT_OK
    assert (target / "design.yaml").is_file()
    assert (target / "layouts.yaml").is_file()
    assert (target / "presentation.yaml").is_file()


def test_init_refuses_nonempty_directory_without_force(tmp_path):
    target = tmp_path / "existing"
    target.mkdir()
    (target / "stray.txt").write_text("hi")
    exit_code = main(["--json", "init", str(target)])
    assert exit_code != EXIT_OK


def test_init_with_no_new_flags_matches_the_bundled_minimal_deck(
    minimal_deck_dir, tmp_path
):
    # Regression guard for issue #34: plain `init` (none of the new
    # --from-dir/--from-repo/--type flags) must stay byte-for-byte
    # identical to the pre-existing bundled-example behavior.
    target = tmp_path / "plain-init"
    exit_code = main(["--json", "init", str(target)])
    assert exit_code == EXIT_OK
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        assert (target / name).read_text() == (minimal_deck_dir / name).read_text()


def _write_flat_template_source(
    root: Path, *, design_name="styleguide.yaml", layouts_name="zones.yaml"
) -> Path:
    root.mkdir(parents=True, exist_ok=True)
    (root / design_name).write_text(yaml.safe_dump({"deckifyr": "0.1", "colors": {}}))
    (root / layouts_name).write_text(yaml.safe_dump({"deckifyr": "0.1", "layouts": {}}))
    (root / "presentation.yaml").write_text(
        yaml.safe_dump(
            {
                "deckifyr": "0.1",
                "design": {"base": design_name},
                "layouts": layouts_name,
                "metadata": {"title": "Source Deck"},
                "build": {"output": "build/source.pptx"},
                "slides": [],
            }
        )
    )
    return root


def test_init_from_dir_flat_scaffolds_a_new_project(tmp_path, capsys):
    source = _write_flat_template_source(tmp_path / "source")
    target = tmp_path / "new-deck"
    exit_code = main(["--json", "init", str(target), "--from-dir", str(source)])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert (target / "styleguide.yaml").is_file()
    assert (target / "zones.yaml").is_file()
    data = yaml.safe_load((target / "presentation.yaml").read_text())
    assert data["slides"] == []
    assert output["warnings"] == []


def test_init_from_dir_typed_requires_type(tmp_path, capsys):
    source = tmp_path / "source"
    _write_flat_template_source(source / "templates" / "alpha")
    exit_code = main(["--json", "init", str(tmp_path / "new-deck"), "--from-dir", str(source)])
    assert exit_code != EXIT_OK
    err = json.loads(capsys.readouterr().err)
    assert "alpha" in err["message"]


def test_init_from_dir_typed_with_type_copies_presentation_verbatim(tmp_path):
    source = tmp_path / "source"
    _write_flat_template_source(source / "templates" / "alpha")
    target = tmp_path / "new-deck"
    exit_code = main(
        ["--json", "init", str(target), "--from-dir", str(source), "--type", "alpha"]
    )
    assert exit_code == EXIT_OK
    data = yaml.safe_load((target / "presentation.yaml").read_text())
    assert data["metadata"]["title"] == "Source Deck"


def test_init_from_dir_unknown_type_lists_available_names(tmp_path, capsys):
    source = tmp_path / "source"
    _write_flat_template_source(source / "templates" / "alpha")
    exit_code = main(
        [
            "--json",
            "init",
            str(tmp_path / "new-deck"),
            "--from-dir",
            str(source),
            "--type",
            "bogus",
        ]
    )
    assert exit_code != EXIT_OK
    err = json.loads(capsys.readouterr().err)
    assert "alpha" in err["message"]


def test_init_from_dir_and_from_repo_are_mutually_exclusive(tmp_path):
    exit_code = main(
        [
            "--json",
            "init",
            str(tmp_path / "new-deck"),
            "--from-dir",
            str(tmp_path),
            "--from-repo",
            "acme/repo",
        ]
    )
    assert exit_code != EXIT_OK


def test_init_ref_without_from_repo_is_an_error(tmp_path):
    exit_code = main(
        ["--json", "init", str(tmp_path / "new-deck"), "--ref", "v1"]
    )
    assert exit_code != EXIT_OK


def test_init_type_without_from_dir_or_from_repo_is_an_error(tmp_path):
    exit_code = main(
        ["--json", "init", str(tmp_path / "new-deck"), "--type", "alpha"]
    )
    assert exit_code != EXIT_OK


def test_init_from_dir_not_a_template_source_is_an_error(tmp_path):
    source = tmp_path / "source"
    source.mkdir()
    exit_code = main(
        ["--json", "init", str(tmp_path / "new-deck"), "--from-dir", str(source)]
    )
    assert exit_code != EXIT_OK


def test_init_from_dir_surfaces_asset_warnings(tmp_path, capsys):
    source = _write_flat_template_source(tmp_path / "source")
    (source / "styleguide.yaml").write_text(
        yaml.safe_dump({"deckifyr": "0.1", "slide": {"background_image": "bg.png"}})
    )
    exit_code = main(
        ["--json", "init", str(tmp_path / "new-deck"), "--from-dir", str(source)]
    )
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert any("bg.png" in w for w in output["warnings"])


def _git(*args: str, cwd: Path) -> None:
    subprocess.run(
        ["git", "-c", "user.email=test@example.com", "-c", "user.name=Test", *args],
        cwd=cwd,
        check=True,
        capture_output=True,
        text=True,
    )


def _make_upstream_repo(root: Path) -> Path:
    repo = root / "upstream"
    repo.mkdir()
    _git("init", cwd=repo)
    _write_flat_template_source(repo)
    _git("add", "-A", cwd=repo)
    _git("commit", "-m", "initial", cwd=repo)
    _git("tag", "v1", cwd=repo)
    return repo


@requires_git
def test_init_from_repo_end_to_end_with_ref_and_subdir(tmp_path):
    repo = tmp_path / "upstream"
    repo.mkdir()
    _git("init", cwd=repo)
    _write_flat_template_source(repo / "templates" / "alpha")
    _git("add", "-A", cwd=repo)
    _git("commit", "-m", "initial", cwd=repo)
    _git("tag", "v1", cwd=repo)

    target = tmp_path / "new-deck"
    exit_code = main(
        [
            "--json",
            "init",
            str(target),
            "--from-repo",
            repo.as_uri(),
            "--ref",
            "v1",
            "--subdir",
            "templates/alpha",
        ]
    )
    assert exit_code == EXIT_OK
    assert (target / "styleguide.yaml").is_file()
    assert (target / "presentation.yaml").is_file()


def test_init_from_repo_missing_git_reports_missing_dependency(tmp_path, monkeypatch, capsys):
    monkeypatch.setattr(shutil, "which", lambda name: None)
    exit_code = main(
        ["--json", "init", str(tmp_path / "new-deck"), "--from-repo", "acme/repo"]
    )
    assert exit_code != EXIT_OK
    err = json.loads(capsys.readouterr().err)
    assert err["code"] == "E_MISSING_DEPENDENCY"
    assert err["dependency"]["name"] == "git"


def test_skills_exports_bundled_skill_files(tmp_path, capsys):
    target = tmp_path / "skills-target"
    exit_code = main(["--json", "skills", str(target)])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    for name in _SKILL_NAMES:
        assert (target / name / "SKILL.md").is_file()
    assert len(output["created"]) == len(_SKILL_NAMES)


def test_skills_refuses_existing_file_without_force(tmp_path):
    target = tmp_path / "skills-target"
    conflicting = target / _SKILL_NAMES[0] / "SKILL.md"
    conflicting.parent.mkdir(parents=True)
    conflicting.write_text("sentinel content")

    exit_code = main(["--json", "skills", str(target)])
    assert exit_code != EXIT_OK
    assert conflicting.read_text() == "sentinel content"


def test_skills_force_overwrites_existing_file(tmp_path):
    target = tmp_path / "skills-target"
    conflicting = target / _SKILL_NAMES[0] / "SKILL.md"
    conflicting.parent.mkdir(parents=True)
    conflicting.write_text("sentinel content")

    exit_code = main(["--json", "skills", str(target), "--force"])
    assert exit_code == EXIT_OK
    assert conflicting.read_text() != "sentinel content"
    assert "name: " + _SKILL_NAMES[0] in conflicting.read_text()


def test_skills_defaults_to_the_current_directory(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    exit_code = main(["--json", "skills"])
    assert exit_code == EXIT_OK
    for name in _SKILL_NAMES:
        assert (tmp_path / name / "SKILL.md").is_file()


def test_inspect_presentation_reports_the_resolved_plan(minimal_deck_dir, capsys):
    exit_code = main(
        ["--json", "inspect", str(minimal_deck_dir / "presentation.yaml")]
    )
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["target"] == "presentation"
    assert output["slide_count"] == 2
    slide_ids = [slide["id"] for slide in output["slides"]]
    assert slide_ids == ["title", "content-slide"]
    assert output["slides"][1]["element_types"] == ["markdown", "text"]


def test_inspect_pptx_reports_real_shape_structure(minimal_deck_dir, tmp_path, capsys):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)
    exit_code = main(["--json", "build", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_OK
    build_output = json.loads(capsys.readouterr().out)

    exit_code = main(["--json", "inspect", build_output["output"]])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["target"] == "pptx"
    assert output["slide_count"] == 2
    shape_names = {
        shape["name"] for slide in output["slides"] for shape in slide["shapes"]
    }
    assert shape_names == {"deck-title", "title", "content"}
    assert output["manifest"]["slide_count"] == 2


def test_inspect_rejects_an_unrecognized_extension(tmp_path, capsys):
    stray = tmp_path / "notes.txt"
    stray.write_text("hi")
    exit_code = main(["--json", "inspect", str(stray)])
    assert exit_code != EXIT_OK
    output = json.loads(capsys.readouterr().err)
    assert output["code"] == "E_IO"


@requires_soffice
def test_preview_renders_one_png_per_slide(minimal_deck_dir, tmp_path, capsys):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)

    exit_code = main(["--json", "preview", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["slide_count"] == 2
    assert len(output["previews"]) == 2
    for preview_path in output["previews"]:
        assert Path(preview_path).is_file()
    # `deckifyr preview` always keeps the intermediate PDF (issue #27) --
    # an ordinary `build.previews: true` build keeps it too now (issue
    # #32), see `test_build_with_previews_enabled_also_keeps_the_pdf`
    # below.
    assert output["preview_pdf"] is not None
    assert Path(output["preview_pdf"]).is_file()


@requires_soffice
def test_preview_slides_flag_renders_only_the_requested_slide(
    minimal_deck_dir, tmp_path, capsys
):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)

    exit_code = main(
        ["--json", "preview", str(tmp_path / "presentation.yaml"), "--slides", "2"]
    )
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["slide_count"] == 2
    assert len(output["previews"]) == 1
    assert Path(output["previews"][0]).name.endswith("-02.png")


def test_preview_slides_flag_rejects_non_integer_input(minimal_deck_dir, tmp_path, capsys):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)

    exit_code = main(
        ["--json", "preview", str(tmp_path / "presentation.yaml"), "--slides", "nope"]
    )
    assert exit_code != EXIT_OK
    output = json.loads(capsys.readouterr().err)
    assert output["code"] == "E_CONTENT_VALIDATION"
