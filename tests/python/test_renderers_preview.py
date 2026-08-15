"""Tests for `deckifyr.renderers.preview` (spec section 12/18 Phase 3).

`_require_soffice` is a pure PATH check -- always runs. `render_slide_previews`
shells out to a real `soffice` binary (and the optional `pymupdf`
dependency), so that test skips cleanly when `soffice` isn't on PATH --
mirroring `test_renderers_quarto.py`'s own quarto-binary skip pattern
(see CLAUDE.md): expected local/CI behavior, not a gap to mock around.
"""

from __future__ import annotations

import shutil

import pytest
from pptx import Presentation

from deckifyr.renderers.preview import PreviewRenderConfig, render_slide_previews
from deckifyr.schema.errors import ContentValidationError, MissingDependencyError

requires_soffice = pytest.mark.skipif(
    shutil.which("soffice") is None, reason="soffice binary not found on PATH"
)


def test_require_soffice_names_the_missing_binary(tmp_path):
    # Exercise the real public entry point rather than the private
    # `_require_soffice` helper directly, so this also proves
    # `render_slide_previews` checks for the binary before touching
    # pymupdf or the filesystem -- `pptx_path`/`out_dir` are deliberately
    # nonsense paths that would fail loudly if reached.
    config = PreviewRenderConfig(binary="definitely-not-a-real-binary")
    with pytest.raises(ContentValidationError, match="definitely-not-a-real-binary"):
        render_slide_previews(
            tmp_path / "does-not-exist.pptx", tmp_path / "out", config=config
        )


def test_require_soffice_raises_a_structured_dependency_error(tmp_path):
    # R/run-python.R's .handle_missing_dependency() reacts to this exact
    # shape (see CLAUDE.md's "Preview rendering" note) -- pin it down so
    # a refactor here can't silently drop the `dependency` payload R
    # depends on.
    config = PreviewRenderConfig(binary="definitely-not-a-real-binary")
    with pytest.raises(MissingDependencyError) as exc_info:
        render_slide_previews(
            tmp_path / "does-not-exist.pptx", tmp_path / "out", config=config
        )
    payload = exc_info.value.to_dict()
    assert payload["code"] == "E_MISSING_DEPENDENCY"
    assert payload["dependency"] == {
        "name": "soffice",
        "display_name": "LibreOffice",
        "install_url": "https://www.libreoffice.org/download/download/",
    }


@requires_soffice
def test_render_slide_previews_produces_one_png_per_slide(tmp_path):
    prs = Presentation()
    blank_layout = next(
        layout for layout in prs.slide_layouts if layout.name == "Blank"
    )
    prs.slides.add_slide(blank_layout)
    prs.slides.add_slide(blank_layout)
    pptx_path = tmp_path / "two-slides.pptx"
    prs.save(str(pptx_path))

    out_dir = tmp_path / "previews"
    image_paths = render_slide_previews(pptx_path, out_dir)

    assert len(image_paths) == 2
    assert [p.name for p in image_paths] == [
        "two-slides-01.png",
        "two-slides-02.png",
    ]
    for image_path in image_paths:
        assert image_path.is_file()
        assert image_path.stat().st_size > 0
