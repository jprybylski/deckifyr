import pathlib
import shutil
import time

import pytest
import yaml
from fastapi.testclient import TestClient
from pptx import Presentation

from deckifyr.schema.presentation import PresentationDocument
from deckifyr.web.app import _frontend_build_warning, _rebuild_frontend, create_app


def _copy_minimal_deck(minimal_deck_dir, tmp_path):
    for name in ("design.yaml", "layouts.yaml", "presentation.yaml"):
        shutil.copyfile(minimal_deck_dir / name, tmp_path / name)
    return tmp_path


@pytest.fixture
def project_dir(minimal_deck_dir, tmp_path):
    return _copy_minimal_deck(minimal_deck_dir, tmp_path)


@pytest.fixture
def client(project_dir):
    app = create_app(project_dir)
    return TestClient(app)


# --- health / project --------------------------------------------------


def test_health(client):
    response = client.get("/api/health")
    assert response.status_code == 200
    body = response.json()
    # `launcher` defaults to "cli" -- see test_health_reports_r_launcher
    # for create_app(..., launcher="r"). `frontend_warning` isn't
    # asserted here -- its real value depends on this checkout's own
    # web/src vs web/static mtimes at test-run time, not anything this
    # test controls; see test_frontend_build_warning_* below for that
    # logic tested in isolation instead.
    assert body["status"] == "ok"
    assert body["launcher"] == "cli"
    assert "frontend_warning" in body


def test_health_reports_r_launcher(tmp_path):
    # /api/health must carry `launcher` even when the bound project fails
    # to load (an empty tmp_path, no presentation.yaml) -- it's the one
    # route the "no project found" screen can rely on to pick CLI- vs
    # R-flavored next-step instructions, precisely when /api/project
    # itself 404s.
    app = create_app(tmp_path, launcher="r")
    with TestClient(app) as local_client:
        response = local_client.get("/api/health")
    assert response.status_code == 200
    assert response.json()["launcher"] == "r"


# --- frontend build-staleness warning (deckifyr.web.app._frontend_build_warning) ---


def test_frontend_build_warning_is_none_outside_a_dev_checkout(tmp_path):
    # No `web/src` sibling at all -- mirrors an installed wheel/R
    # package, which never ships the frontend source tree.
    assert _frontend_build_warning(tmp_path / "static", tmp_path / "web" / "src") is None


def test_frontend_build_warning_is_none_when_the_build_is_current(tmp_path):
    web_src = tmp_path / "web" / "src"
    static_dir = tmp_path / "static"
    web_src.mkdir(parents=True)
    static_dir.mkdir(parents=True)
    (web_src / "App.tsx").write_text("// source")
    time.sleep(0.01)
    (static_dir / "index.html").write_text("<html></html>")
    assert _frontend_build_warning(static_dir, web_src) is None


def test_frontend_build_warning_fires_when_source_is_newer_than_the_build(tmp_path):
    web_src = tmp_path / "web" / "src"
    static_dir = tmp_path / "static"
    web_src.mkdir(parents=True)
    static_dir.mkdir(parents=True)
    (static_dir / "index.html").write_text("<html></html>")
    time.sleep(0.01)
    (web_src / "App.tsx").write_text("// a real fix, not yet compiled")
    warning = _frontend_build_warning(static_dir, web_src)
    assert warning is not None
    assert "npm run build" in warning


def test_frontend_build_warning_ignores_test_files(tmp_path):
    # A touched *.test.tsx must not itself trigger the warning -- only
    # real source files count, per the module's own docstring.
    web_src = tmp_path / "web" / "src"
    static_dir = tmp_path / "static"
    web_src.mkdir(parents=True)
    static_dir.mkdir(parents=True)
    (web_src / "App.tsx").write_text("// source")
    time.sleep(0.01)
    (static_dir / "index.html").write_text("<html></html>")
    time.sleep(0.01)
    (web_src / "App.test.tsx").write_text("// a newer test-only edit")
    assert _frontend_build_warning(static_dir, web_src) is None


def test_frontend_build_warning_fires_when_no_build_exists_yet(tmp_path):
    web_src = tmp_path / "web" / "src"
    web_src.mkdir(parents=True)
    (web_src / "App.tsx").write_text("// source")
    warning = _frontend_build_warning(tmp_path / "static", web_src)
    assert warning is not None


# --- frontend auto-rebuild (deckifyr.web.app._rebuild_frontend) ---


def test_rebuild_frontend_is_a_silent_noop_without_npm_on_path(tmp_path, monkeypatch):
    monkeypatch.setattr(shutil, "which", lambda name: None)
    assert _rebuild_frontend(tmp_path) is None


def test_rebuild_frontend_reports_a_failed_build(tmp_path, monkeypatch):
    import subprocess as subprocess_module

    monkeypatch.setattr(shutil, "which", lambda name: "/usr/bin/npm")

    def fake_run(cmd, cwd, capture_output, text, timeout):
        return subprocess_module.CompletedProcess(cmd, returncode=1, stdout="", stderr="a real build error\n")

    monkeypatch.setattr("deckifyr.web.app.subprocess.run", fake_run)
    message = _rebuild_frontend(tmp_path)
    assert message is not None
    assert "npm run build failed" in message
    assert "a real build error" in message


def test_rebuild_frontend_returns_none_on_a_successful_build(tmp_path, monkeypatch):
    import subprocess as subprocess_module

    monkeypatch.setattr(shutil, "which", lambda name: "/usr/bin/npm")

    def fake_run(cmd, cwd, capture_output, text, timeout):
        return subprocess_module.CompletedProcess(cmd, returncode=0, stdout="", stderr="")

    monkeypatch.setattr("deckifyr.web.app.subprocess.run", fake_run)
    assert _rebuild_frontend(tmp_path) is None


@pytest.mark.skipif(shutil.which("npm") is None, reason="npm not on PATH")
def test_rebuild_frontend_really_rebuilds_the_repo_frontend():
    # End-to-end against this repo's real web/ checkout (already has
    # node_modules from earlier `npm install`/`npm test` runs in CI/dev --
    # see tests/python/test_renderers_quarto.py's own npm-equivalent
    # skip-when-absent precedent for real-toolchain integration tests).
    web_src_dir = pathlib.Path(__file__).resolve().parents[2] / "web" / "src"
    if not (web_src_dir.parent / "node_modules").is_dir():
        pytest.skip("web/node_modules not installed -- run npm install in web/ first")
    static_dir = pathlib.Path(__file__).resolve().parents[2] / "inst" / "python" / "deckifyr" / "web" / "static"
    assert _rebuild_frontend(web_src_dir.parent) is None
    assert _frontend_build_warning(static_dir, web_src_dir) is None


def test_project_reports_bound_paths(client, project_dir):
    response = client.get("/api/project")
    assert response.status_code == 200
    body = response.json()
    assert body["root"] == str(project_dir.resolve())
    assert body["presentation"].endswith("presentation.yaml")
    assert body["design"].endswith("design.yaml")
    assert body["layouts"].endswith("layouts.yaml")


# --- config get/put ------------------------------------------------------


def test_get_config_design(client):
    response = client.get("/api/config/design")
    assert response.status_code == 200
    assert response.json()["colors"]["primary"] == "#2457A6"


def test_get_config_unknown_doc_is_404(client):
    response = client.get("/api/config/nope")
    assert response.status_code == 404


def test_put_config_design_valid_edit_round_trips(client, project_dir):
    design_path = project_dir / "design.yaml"
    original = design_path.read_text()
    data = yaml.safe_load(original)
    data["colors"]["primary"] = "#123456"

    response = client.put("/api/config/design", json=data)
    assert response.status_code == 200
    assert response.json()["dirty"] is True

    # Deferred-save editing (issue #24): an edit lands in the in-memory
    # working copy, not the file, until an explicit `POST /api/save` --
    # the whole point is that trying something in the web editor is safe
    # without dirtying a tracked project's YAML.
    assert design_path.read_text() == original

    response = client.get("/api/config/design")
    assert response.json()["colors"]["primary"] == "#123456"


def test_put_config_design_invalid_edit_is_422_and_file_unchanged(client, project_dir):
    design_path = project_dir / "design.yaml"
    original = design_path.read_text()
    data = yaml.safe_load(original)
    # `slide.width` is required -- dropping it breaks schema validation.
    del data["slide"]["width"]

    response = client.put("/api/config/design", json=data)
    assert response.status_code == 422
    body = response.json()
    assert "code" in body and "message" in body

    assert design_path.read_text() == original


def test_put_config_presentation_invalid_layout_reference_is_422(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    original = presentation_path.read_text()
    data = yaml.safe_load(original)
    data["slides"][1]["layout"] = "does-not-exist"

    response = client.put("/api/config/presentation", json=data)
    assert response.status_code == 422
    assert presentation_path.read_text() == original


# --- plan ------------------------------------------------------------------


def test_plan_has_expected_slide_and_element_shape_with_formatted_geometry(client):
    response = client.get("/api/plan")
    assert response.status_code == 200
    body = response.json()

    slides = body["slides"]
    assert len(slides) == 2
    assert slides[0]["id"] == "title"

    elements = slides[0]["elements"]
    assert any(el["id"] == "deck-title" for el in elements)
    deck_title = next(el for el in elements if el["id"] == "deck-title")

    # Geometry is formatted through `format_length` as unit strings, not
    # raw EMU ints -- the box the minimal deck fixture declares is
    # `{x: 0.9in, y: 2.1in, width: 11.5in, height: 1.1in}`.
    assert deck_title["box"] == {
        "x": "0.9in",
        "y": "2.1in",
        "width": "11.5in",
        "height": "1.1in",
    }
    assert "x" not in deck_title
    assert "width" not in deck_title


# --- element PATCH -----------------------------------------------------


def test_patch_element_box_and_rotation_updates_working_copy_and_response(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    original = presentation_path.read_text()

    response = client.patch(
        "/api/slides/title/elements/deck-title",
        json={"box": {"x": 1.0, "y": 2.5}, "rotation": 15},
    )
    assert response.status_code == 200
    assert response.json()["dirty"] is True

    # Not written to disk until Save -- see test_put_config_design_valid_edit_round_trips.
    assert presentation_path.read_text() == original

    plan_response = client.get("/api/plan")
    elements = plan_response.json()["slides"][0]["elements"]
    element = next(el for el in elements if el["id"] == "deck-title")
    # `GET /api/plan` formats geometry through `format_length` (spec
    # section 7.3), which normalizes "1.0in" -> "1in" -- unlike the old
    # version of this test, which read the raw YAML string back verbatim.
    assert element["box"]["x"] == "1in"
    assert element["box"]["y"] == "2.5in"
    assert element["rotation"] == 15


def test_patch_element_unknown_slide_is_404(client):
    response = client.patch(
        "/api/slides/does-not-exist/elements/deck-title", json={"rotation": 5}
    )
    assert response.status_code == 404


def test_patch_element_unknown_element_is_404(client):
    response = client.patch(
        "/api/slides/title/elements/does-not-exist", json={"rotation": 5}
    )
    assert response.status_code == 404


def test_patch_element_invalid_value_is_422_and_file_unchanged(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    original = presentation_path.read_text()

    # `rotation` is a float field -- a non-numeric value fails schema
    # validation (unlike a negative box dimension, which `Box`'s plain
    # `str` fields don't reject at this layer at all).
    response = client.patch(
        "/api/slides/title/elements/deck-title", json={"rotation": "sideways"}
    )
    assert response.status_code == 422

    assert presentation_path.read_text() == original


# --- slide add/remove (issue #23) ---------------------------------------


def test_add_slide_appends_by_default(client, project_dir):
    response = client.post("/api/slides", json={"id": "new-slide", "layout": "blank"})
    assert response.status_code == 200
    body = response.json()
    assert body["slide_count"] == 3
    assert body["dirty"] is True

    plan = client.get("/api/plan").json()
    assert [s["id"] for s in plan["slides"]] == ["title", "content-slide", "new-slide"]
    assert plan["slide_layouts"]["new-slide"] == "blank"

    # Not written to disk until Save.
    presentation_path = project_dir / "presentation.yaml"
    assert "new-slide" not in presentation_path.read_text()


def test_add_slide_with_freeform_layout_and_placement(client):
    response = client.post(
        "/api/slides", json={"id": "freeform", "layout": None, "after": "title"}
    )
    assert response.status_code == 200
    plan = client.get("/api/plan").json()
    assert [s["id"] for s in plan["slides"]] == ["title", "freeform", "content-slide"]
    assert plan["slide_layouts"]["freeform"] is None


def test_add_slide_duplicate_id_is_422(client):
    response = client.post("/api/slides", json={"id": "title", "layout": "blank"})
    assert response.status_code == 422


def test_add_slide_unknown_after_reference_is_422(client):
    response = client.post(
        "/api/slides", json={"id": "new-slide", "layout": "blank", "after": "nope"}
    )
    assert response.status_code == 422


def test_add_slide_ambiguous_placement_is_422(client):
    response = client.post(
        "/api/slides",
        json={"id": "new-slide", "layout": "blank", "after": "title", "index": 0},
    )
    assert response.status_code == 422


def test_remove_slide_updates_working_copy(client, project_dir):
    response = client.delete("/api/slides/content-slide")
    assert response.status_code == 200
    body = response.json()
    assert body["slide_count"] == 1
    assert body["dirty"] is True

    plan = client.get("/api/plan").json()
    assert [s["id"] for s in plan["slides"]] == ["title"]

    presentation_path = project_dir / "presentation.yaml"
    assert "content-slide" in presentation_path.read_text()  # unchanged until Save


def test_remove_slide_unknown_id_is_422(client):
    response = client.delete("/api/slides/does-not-exist")
    assert response.status_code == 422


# --- layout zones (issue #23's Content/Layout tab) -----------------------


def test_get_layout_zones_returns_resolved_boxes(client):
    response = client.get("/api/layouts/title-content")
    assert response.status_code == 200
    body = response.json()
    assert body["id"] == "__layout__title-content"
    zone_ids = {el["id"] for el in body["elements"]}
    assert zone_ids == {"title", "content", "footnotes"}
    title_zone = next(el for el in body["elements"] if el["id"] == "title")
    assert title_zone["box"] == {
        "x": "0.7in",
        "y": "0.35in",
        "width": "11.9in",
        "height": "0.65in",
    }


def test_get_layout_zones_empty_layout_has_no_elements(client):
    response = client.get("/api/layouts/blank")
    assert response.status_code == 200
    assert response.json()["elements"] == []


def test_get_layout_zones_unknown_layout_is_404(client):
    response = client.get("/api/layouts/does-not-exist")
    assert response.status_code == 404


def test_patch_layout_element_updates_box_and_working_copy(client, project_dir):
    layouts_path = project_dir / "layouts.yaml"
    original = layouts_path.read_text()

    response = client.patch(
        "/api/layouts/title-content/elements/title",
        json={"box": {"x": 1.0, "y": 0.5}},
    )
    assert response.status_code == 200
    assert response.json()["dirty"] is True
    assert layouts_path.read_text() == original  # not written until Save

    body = client.get("/api/layouts/title-content").json()
    title_zone = next(el for el in body["elements"] if el["id"] == "title")
    assert title_zone["box"]["x"] == "1in"
    assert title_zone["box"]["y"] == "0.5in"

    # Every slide using this layout sees the change too (shared state).
    plan = client.get("/api/plan").json()
    content_slide = next(s for s in plan["slides"] if s["id"] == "content-slide")
    title_el = next(el for el in content_slide["elements"] if el["id"] == "title")
    assert title_el["box"]["x"] == "1in"


def test_patch_layout_element_unknown_zone_is_422(client):
    response = client.patch(
        "/api/layouts/title-content/elements/does-not-exist",
        json={"box": {"x": 1, "y": 1}},
    )
    assert response.status_code == 422


# --- furniture pseudo-slide (issue #21) ---------------------------------


def test_get_furniture_empty_when_nothing_configured(client):
    response = client.get("/api/furniture")
    assert response.status_code == 200
    body = response.json()
    assert body["id"] == "__furniture__"
    assert body["elements"] == []


def test_patch_furniture_element_not_configured_is_422(client):
    response = client.patch(
        "/api/furniture/elements/__furniture_branding",
        json={"box": {"x": 1, "y": 1, "width": 2, "height": 0.5}},
    )
    assert response.status_code == 422
    assert response.json()["code"] == "E_PATH_NOT_FOUND"


def test_patch_furniture_element_background_is_422(client):
    response = client.patch(
        "/api/furniture/elements/__furniture_background",
        json={"box": {"x": 0, "y": 0, "width": 1, "height": 1}},
    )
    assert response.status_code == 422


def test_add_patch_and_remove_branding_furniture(client, project_dir):
    design_path = project_dir / "design.yaml"
    original = design_path.read_text()

    add_response = client.post("/api/furniture/elements/__furniture_branding")
    assert add_response.status_code == 200
    assert add_response.json()["dirty"] is True
    # Not written to disk until Save -- see test_put_config_design_valid_edit_round_trips.
    assert design_path.read_text() == original
    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["branding"] is not None

    # A second add is a conflict, not a silent overwrite.
    conflict_response = client.post("/api/furniture/elements/__furniture_branding")
    assert conflict_response.status_code == 422

    plan_response = client.get("/api/furniture")
    elements = plan_response.json()["elements"]
    assert any(el["id"] == "__furniture_branding" for el in elements)

    patch_response = client.patch(
        "/api/furniture/elements/__furniture_branding",
        json={"box": {"x": 1.0, "y": 2.0, "width": 3.0, "height": 0.5}, "value": "Acme Corp"},
    )
    assert patch_response.status_code == 200
    on_disk = client.get("/api/config/design").json()
    branding = on_disk["furniture"]["branding"]
    assert branding["box"]["x"] == "1.0in"
    assert branding["text"] == "Acme Corp"

    # branding has no rotation/z_index field in the schema -- either is a
    # hard reject, not a silent no-op.
    rotation_response = client.patch(
        "/api/furniture/elements/__furniture_branding", json={"rotation": 10}
    )
    assert rotation_response.status_code == 422
    z_index_response = client.patch(
        "/api/furniture/elements/__furniture_branding", json={"z_index": 5}
    )
    assert z_index_response.status_code == 422

    remove_response = client.delete("/api/furniture/elements/__furniture_branding")
    assert remove_response.status_code == 200
    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["branding"] is None

    plan_response = client.get("/api/furniture")
    assert plan_response.json()["elements"] == []


def test_add_and_patch_page_number_furniture(client, project_dir):
    design_path = project_dir / "design.yaml"

    add_response = client.post("/api/furniture/elements/__furniture_page_number")
    assert add_response.status_code == 200

    patch_response = client.patch(
        "/api/furniture/elements/__furniture_page_number",
        json={"box": {"x": 0.5, "y": 0.5, "width": 1.0, "height": 0.3}},
    )
    assert patch_response.status_code == 200
    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["page_number"]["box"]["x"] == "0.5in"

    # page_number.format isn't editable through `value` -- Config-tab only.
    value_response = client.patch(
        "/api/furniture/elements/__furniture_page_number", json={"value": "p. {page}"}
    )
    assert value_response.status_code == 422


def test_add_furniture_status_defaults_to_watermark_when_nothing_selected(client, project_dir):
    # "Add" is now the quick, one-click way to get *a* status indicator
    # at all -- with nothing selected yet (status_indicator unset), it
    # defaults to the watermark placement rather than requiring the Deck
    # Options dropdown to be touched first (issue #24's follow-up revert:
    # see `deckifyr.plan.FURNITURE_STATUS_ID`'s own docstring for why the
    # separate, additive `watermark_overlay` concept this replaced was
    # dropped). minimal-deck's own `metadata.status` ("draft") supplies
    # the text.
    add_response = client.post("/api/furniture/elements/__furniture_status")
    assert add_response.status_code == 200

    presentation_after = client.get("/api/config/presentation").json()
    assert presentation_after["status_indicator"] == "watermark"

    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["status"]["watermark"] is not None

    plan_response = client.get("/api/furniture")
    status_element = next(
        el for el in plan_response.json()["elements"] if el["id"] == "__furniture_status"
    )
    assert status_element["value"] == "draft"


def test_add_furniture_status_defaulting_to_watermark_fails_with_no_text_source(
    client, project_dir
):
    presentation_path = project_dir / "presentation.yaml"
    data = yaml.safe_load(presentation_path.read_text())
    data["metadata"]["status"] = None
    client.put("/api/config/presentation", json=data)

    # Same validation `_check_watermark_has_text` always enforced for a
    # `status_indicator: watermark` selected via the dropdown -- "Add"
    # defaulting to it is not a way around that check.
    add_response = client.post("/api/furniture/elements/__furniture_status")
    assert add_response.status_code == 422
    assert client.get("/api/config/presentation").json().get("status_indicator") is None


def test_status_furniture_corner_add_patch_and_remove(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    data = yaml.safe_load(presentation_path.read_text())
    data["status_indicator"] = "corner-br"
    # A corner placement always shows metadata.status ("draft" in the
    # minimal-deck fixture), never `watermark` -- that override only
    # applies to the full-page `"watermark"` placement (see
    # `resolve_watermark_text`'s own docstring). Set here specifically to
    # confirm it's ignored for a corner, not used to derive the assertion.
    data["watermark"] = "DRAFT"
    put_response = client.put("/api/config/presentation", json=data)
    assert put_response.status_code == 200

    # A corner already selected via the dropdown -- "Add" here just
    # materializes its style, doesn't touch status_indicator (already
    # set to the right thing).
    add_response = client.post("/api/furniture/elements/__furniture_status")
    assert add_response.status_code == 200
    assert client.get("/api/config/presentation").json()["status_indicator"] == "corner-br"

    plan_response = client.get("/api/furniture")
    status_element = next(
        el for el in plan_response.json()["elements"] if el["id"] == "__furniture_status"
    )
    assert status_element["value"] == "draft"

    patch_response = client.patch(
        "/api/furniture/elements/__furniture_status",
        json={"rotation": 12, "z_index": 250},
    )
    assert patch_response.status_code == 200
    on_disk = client.get("/api/config/design").json()
    corner_br = on_disk["furniture"]["status"]["corner_br"]
    assert corner_br["rotation"] == 12
    assert corner_br["z_index"] == 250

    remove_response = client.delete("/api/furniture/elements/__furniture_status")
    assert remove_response.status_code == 200
    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["status"]["corner_br"] is None

    # Removing the active placement also clears presentation.yaml's
    # status_indicator in the same action -- leaving it pointing at a
    # now-nonexistent style would be exactly the dangling-reference
    # footgun Remove used to be withheld to avoid. The override *text*
    # (`watermark`) is a content choice, not an activation flag, so it's
    # left alone.
    presentation_after = client.get("/api/config/presentation").json()
    assert presentation_after["status_indicator"] is None
    assert presentation_after["watermark"] == "DRAFT"


def test_remove_furniture_status_while_watermark_is_active(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    data = yaml.safe_load(presentation_path.read_text())
    data["status_indicator"] = "watermark"
    data["watermark"] = "test"
    client.put("/api/config/presentation", json=data)
    assert client.post("/api/furniture/elements/__furniture_status").status_code == 200

    remove_response = client.delete("/api/furniture/elements/__furniture_status")
    assert remove_response.status_code == 200

    on_disk = client.get("/api/config/design").json()
    assert on_disk["furniture"]["status"]["watermark"] is None
    presentation_after = client.get("/api/config/presentation").json()
    assert presentation_after["status_indicator"] is None
    # The override *text* is a content choice, not an activation flag --
    # stays useful if the watermark is selected again later.
    assert presentation_after["watermark"] == "test"


def test_get_furniture_is_lenient_when_status_indicator_points_at_an_unconfigured_placement(
    client, project_dir
):
    # Regression: a real user hit this exact scenario -- picking a new
    # status_indicator placement in Deck Options before design.yaml has
    # a style for it. Before this fix, GET /api/furniture raised the
    # same ContentValidationError GET /api/plan does, which left
    # FurnitureControls with nothing to render and no way to reach its
    # own "Add" fix -- the one screen meant to recover from this got
    # stuck too. `furniture.branding` is configured first specifically
    # to confirm a failing status placement doesn't take unrelated,
    # already-configured furniture down with it.
    add_response = client.post("/api/furniture/elements/__furniture_branding")
    assert add_response.status_code == 200

    presentation_path = project_dir / "presentation.yaml"
    data = yaml.safe_load(presentation_path.read_text())
    data["status_indicator"] = "corner-tl"
    data["watermark"] = "DRAFT"
    put_response = client.put("/api/config/presentation", json=data)
    assert put_response.status_code == 200

    # GET /api/plan (real-slide rendering) stays strict -- unchanged.
    plan_response = client.get("/api/plan")
    assert plan_response.status_code == 422

    # GET /api/furniture degrades gracefully: 200, __furniture_status is
    # simply absent, and __furniture_branding still comes through.
    furniture_response = client.get("/api/furniture")
    assert furniture_response.status_code == 200
    element_ids = {el["id"] for el in furniture_response.json()["elements"]}
    assert "__furniture_status" not in element_ids
    assert "__furniture_branding" in element_ids


def test_furniture_element_unknown_id_is_404(client):
    assert client.patch("/api/furniture/elements/nope", json={}).status_code == 404
    assert client.post("/api/furniture/elements/nope").status_code == 404
    assert client.delete("/api/furniture/elements/nope").status_code == 404


# --- save / discard / autosave (issue #24) -------------------------------


def test_save_writes_only_touched_documents(client, project_dir):
    design_path = project_dir / "design.yaml"
    presentation_path = project_dir / "presentation.yaml"
    design_before = design_path.read_text()
    presentation_before = presentation_path.read_text()

    data = yaml.safe_load(design_before)
    data["colors"]["primary"] = "#654321"
    response = client.put("/api/config/design", json=data)
    assert response.json()["dirty"] is True

    # Not written to disk until Save.
    assert design_path.read_text() == design_before
    assert presentation_path.read_text() == presentation_before

    save_response = client.post("/api/save")
    assert save_response.status_code == 200
    assert save_response.json() == {"saved": ["design"], "dirty": False}

    # Only the touched document was written -- presentation.yaml, never
    # edited this session, keeps its original bytes untouched.
    assert design_path.read_text() != design_before
    assert presentation_path.read_text() == presentation_before
    assert yaml.safe_load(design_path.read_text())["colors"]["primary"] == "#654321"


def test_discard_reverts_unsaved_edits(client, project_dir):
    design_path = project_dir / "design.yaml"
    original = design_path.read_text()

    data = yaml.safe_load(original)
    data["colors"]["primary"] = "#abcdef"
    response = client.put("/api/config/design", json=data)
    assert response.json()["dirty"] is True

    discard_response = client.post("/api/discard")
    assert discard_response.status_code == 200
    assert discard_response.json() == {"dirty": False}

    # The working copy re-reads from disk on next access -- back to the
    # original value, and the file itself was never touched either.
    assert client.get("/api/config/design").json()["colors"]["primary"] != "#abcdef"
    assert design_path.read_text() == original


def test_autosave_flushes_immediately_with_no_separate_save_call(client, project_dir):
    presentation_path = project_dir / "presentation.yaml"
    design_path = project_dir / "design.yaml"

    data = yaml.safe_load(presentation_path.read_text())
    data.setdefault("build", {})["autosave"] = True
    response = client.put("/api/config/presentation", json=data)
    assert response.status_code == 200
    # Turning autosave on is itself an edit, and it flushes right away --
    # no special-casing needed, since the mutation that sets `autosave:
    # true` is itself what `_after_mutation` reacts to.
    assert response.json()["dirty"] is False
    assert yaml.safe_load(presentation_path.read_text())["build"]["autosave"] is True

    design_data = yaml.safe_load(design_path.read_text())
    design_data["colors"]["primary"] = "#00ff00"
    design_response = client.put("/api/config/design", json=design_data)
    assert design_response.status_code == 200
    assert design_response.json()["dirty"] is False
    assert yaml.safe_load(design_path.read_text())["colors"]["primary"] == "#00ff00"


def test_get_plan_reports_dirty_state(client, project_dir):
    assert client.get("/api/plan").json()["dirty"] is False

    data = yaml.safe_load((project_dir / "presentation.yaml").read_text())
    client.put("/api/config/presentation", json=data)
    assert client.get("/api/plan").json()["dirty"] is True

    client.post("/api/save")
    assert client.get("/api/plan").json()["dirty"] is False


# --- build job lifecycle ----------------------------------------------


def test_build_job_lifecycle(client):
    response = client.post("/api/build")
    assert response.status_code == 200
    job_id = response.json()["job_id"]

    status = None
    deadline = time.monotonic() + 30
    while time.monotonic() < deadline:
        job_response = client.get(f"/api/jobs/{job_id}")
        assert job_response.status_code == 200
        status = job_response.json()["status"]
        if status in ("succeeded", "failed"):
            break
        time.sleep(0.25)

    job_body = job_response.json()
    assert status == "succeeded", job_body
    assert job_body["result"] is not None
    assert job_body["error"] is None

    artifacts_response = client.get(f"/api/jobs/{job_id}/artifacts")
    assert artifacts_response.status_code == 200
    artifact_keys = artifacts_response.json()["artifacts"]
    assert "pptx" in artifact_keys

    pptx_response = client.get(f"/api/jobs/{job_id}/artifacts/pptx")
    assert pptx_response.status_code == 200
    assert len(pptx_response.content) > 0
    # A real, readable .pptx package -- not just nonzero bytes.
    import io

    Presentation(io.BytesIO(pptx_response.content))

    missing_response = client.get(f"/api/jobs/{job_id}/artifacts/nonexistent-key")
    assert missing_response.status_code == 404


def test_job_not_found_is_404(client):
    response = client.get("/api/jobs/does-not-exist")
    assert response.status_code == 404


# --- preview (issue #27) -------------------------------------------------

requires_soffice = pytest.mark.skipif(
    shutil.which("soffice") is None, reason="soffice binary not found on PATH"
)


def test_preview_availability_reports_soffice_on_path(client):
    response = client.get("/api/preview/availability")
    assert response.status_code == 200
    body = response.json()
    assert body["binary"] == "soffice"
    assert body["available"] is (shutil.which("soffice") is not None)
    if body["available"]:
        assert body["install_url"] is None
    else:
        assert body["install_url"] == "https://www.libreoffice.org/download/download/"


def test_preview_availability_reports_unavailable_for_a_bogus_binary(
    client, project_dir
):
    doc = yaml.safe_load((project_dir / "presentation.yaml").read_text())
    doc["build"]["preview"] = {"binary": "definitely-not-a-real-binary"}
    put_response = client.put("/api/config/presentation", json=doc)
    assert put_response.status_code == 200

    response = client.get("/api/preview/availability")
    assert response.status_code == 200
    body = response.json()
    assert body["binary"] == "definitely-not-a-real-binary"
    assert body["available"] is False
    assert body["display_name"] == "LibreOffice"
    assert body["install_url"] == "https://www.libreoffice.org/download/download/"


@requires_soffice
def test_preview_job_lifecycle_produces_selected_previews_and_pdf(client):
    response = client.post("/api/preview", json={"slides": [1]})
    assert response.status_code == 200
    job_id = response.json()["job_id"]

    status = None
    deadline = time.monotonic() + 60
    while time.monotonic() < deadline:
        job_response = client.get(f"/api/jobs/{job_id}")
        assert job_response.status_code == 200
        status = job_response.json()["status"]
        if status in ("succeeded", "failed"):
            break
        time.sleep(0.25)

    job_body = job_response.json()
    assert status == "succeeded", job_body

    artifact_keys = client.get(f"/api/jobs/{job_id}/artifacts").json()["artifacts"]
    assert set(artifact_keys) == {"pdf", "preview-0", "pptx"}
    assert "preview-1" not in artifact_keys  # only slide 1 was requested

    pdf_response = client.get(f"/api/jobs/{job_id}/artifacts/pdf")
    assert pdf_response.status_code == 200
    assert pdf_response.content[:4] == b"%PDF"


def test_post_preview_with_no_body_defaults_to_all_slides(client):
    response = client.post("/api/preview")
    assert response.status_code == 200
    assert "job_id" in response.json()


# --- schemas -------------------------------------------------------------


def test_schema_presentation_matches_model(client):
    response = client.get("/api/schemas/presentation")
    assert response.status_code == 200
    assert response.json() == PresentationDocument.model_json_schema()


def test_schema_unknown_doc_is_404(client):
    response = client.get("/api/schemas/nope")
    assert response.status_code == 404
