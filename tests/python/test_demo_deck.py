"""End-to-end coverage for examples/demo-deck (see its README.md): a
six-slide deck built from a real reportifyr-produced figure (resolved
via a real `{rpfy}:` magic string, spec section 9), a derived CSV table,
a `quarto`-fragment slide (spec section 8.1, issue #3), and a
reportifyr-table-formats slide (issue #57) -- not placeholder content.
This is a regression test for the whole plan -> compose pipeline against
a richer project than inst/examples/minimal-deck's text/markdown-only
fixture -- in particular, it's the only test that builds a project with
a multi-zone layout, `rotation`, `z_index`, a `table` element, real
`quarto` elements (one rendered `svg`/`png`-style as `render_mode: png`,
one executing a real R code chunk as `render_mode: native`), and a
`.rds` flextable reportifyr artifact all at once.

The whole deck now requires a real `quarto` binary to build (the
`pk-interpretation` slide's two elements), so every test in this file
skips cleanly when `quarto` isn't on `PATH` -- and the R-executing
fragment additionally needs `Rscript`, and the `table-formats` slide's
`.rds` flextable artifact additionally needs the R `flextable` package
on top of that (`deckifyr.renderers.flextable` shells to `Rscript`
directly for this slide, not just Quarto executing an R chunk) --
mirroring `tests/python/test_renderers_quarto.py`'s own skip pattern
(see CLAUDE.md's "Quarto integration" architecture note). This is
expected local/CI behavior: CI's `python-tests` job does not install
Quarto, R, or the `flextable` R package, so this file's real pipeline
coverage runs wherever all three are available (this was run and
verified locally against a live `quarto`/R/`flextable` install) rather
than showing up in CI-tracked coverage numbers.
"""

from __future__ import annotations

import hashlib
import json
import shutil
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from deckifyr.cli import EXIT_OK, main

requires_quarto = pytest.mark.skipif(
    shutil.which("quarto") is None, reason="quarto binary not found on PATH"
)
requires_r = pytest.mark.skipif(
    shutil.which("Rscript") is None, reason="Rscript not found on PATH"
)


def _rscript_has_flextable() -> bool:
    if shutil.which("Rscript") is None:
        return False
    try:
        result = subprocess.run(
            [
                "Rscript",
                "--vanilla",
                "-e",
                "quit(status = if (requireNamespace('flextable', quietly = TRUE)) 0 else 1)",
            ],
            capture_output=True,
            timeout=30,
        )
    except (OSError, subprocess.TimeoutExpired):
        return False
    return result.returncode == 0


requires_flextable = pytest.mark.skipif(
    not _rscript_has_flextable(), reason="Rscript with the flextable package not found"
)
pytestmark = [requires_quarto, requires_r, requires_flextable]


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _build_demo_deck(demo_deck_dir: Path, tmp_path: Path, capsys) -> dict:
    # Copy the whole project (including OUTPUTS/, fragments/, and
    # assets/, not just the three YAML files) into tmp_path so the
    # build's output doesn't land inside the repo checkout.
    shutil.copytree(demo_deck_dir, tmp_path, dirs_exist_ok=True, ignore=shutil.ignore_patterns("build"))

    exit_code = main(["--json", "build", str(tmp_path / "presentation.yaml")])
    assert exit_code == EXIT_OK
    output = json.loads(capsys.readouterr().out)
    assert output["status"] == "ok"
    return output


def _alt_text(picture_shape) -> str | None:
    cnv_pr = picture_shape._element.find(qn("p:nvPicPr")).find(qn("p:cNvPr"))
    return cnv_pr.get("descr")


def test_demo_deck_builds_six_slides_with_expected_shapes(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    assert output["slide_count"] == 6

    prs = Presentation(output["output"])
    slides = list(prs.slides)
    assert len(slides) == 6

    furniture = {
        "__furniture_background",
        "__furniture_status",
        "__furniture_branding",
        "__furniture_page_number",
    }
    shape_names_per_slide = [{shape.name for shape in slide.shapes} for slide in slides]
    assert shape_names_per_slide == [
        {"deck-title", "deck-subtitle"} | furniture,
        {"title", "figure", "figure__footer", "note"} | furniture,
        {"title", "figure", "note"} | furniture,
        {"table-title", "pk-table"} | furniture,
        {
            "table-formats-title",
            "raw-table",
            "raw-table__footer",
            "flextable-summary",
            "flextable-summary__footer",
        }
        | furniture,
        {"closing-title", "closing-note", "logo"} | furniture,
    ]


def test_demo_deck_figure_is_a_picture_with_alt_text(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    slides = list(prs.slides)
    plot_slide = slides[1]

    figure = next(shape for shape in plot_slide.shapes if shape.name == "figure")
    assert figure.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    assert "Theoph" in _alt_text(figure)


def test_demo_deck_figure_footer_reflects_the_real_metadata_sidecar(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    slides = list(prs.slides)
    plot_slide = slides[1]

    footer = next(shape for shape in plot_slide.shapes if shape.name == "figure__footer")
    footer_text = footer.text_frame.text
    assert footer_text.startswith("Source: scripts/01_analysis.R")
    assert "Theoph dataset" in footer_text
    assert "Abbreviations: PK: pharmacokinetic." in footer_text


def test_demo_deck_pk_table_is_a_native_table_from_csv(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    slides = list(prs.slides)
    pk_slide = slides[3]

    graphic_frame = next(shape for shape in pk_slide.shapes if shape.name == "pk-table")
    table = graphic_frame.table
    assert [cell.text for cell in table.rows[0].cells] == [
        "Participant", "Weight (kg)", "Dose (mg/kg)", "Cmax (mg/L)", "Tmax (hr)",
    ]
    # 12 Theoph participants, one data row each, plus the header row.
    assert len(table.rows) == 13
    assert [cell.text for cell in table.rows[1].cells] == ["1", "79.6", "4.02", "10.50", "1.12"]


def test_demo_deck_logo_keeps_its_rotation(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    closing_slide = list(prs.slides)[5]

    logo = next(shape for shape in closing_slide.shapes if shape.name == "logo")
    # presentation.yaml sets rotation: -3; python-pptx normalizes negative
    # rotation to [0, 360), so -3 degrees clockwise reads back as 357.0.
    assert logo.rotation == 357.0
    assert _alt_text(logo) == "Organization logo"


def test_demo_deck_plot_slide_carries_its_speaker_notes(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    slides = list(prs.slides)
    plot_slide = slides[1]
    closing_slide = slides[5]

    assert "absorption phase" in plot_slide.notes_slide.notes_text_frame.text
    assert closing_slide.has_notes_slide is False


def test_demo_deck_manifest_records_the_real_figure_hash(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    manifest = json.loads(Path(output["manifest"]).read_text())

    assert manifest["slide_count"] == 6
    # 2 + 3 + 3 + 2 + 3 + 3 elements across the six slides (title/
    # deck-subtitle; title/figure/note; title/figure/note; table-title/
    # pk-table; table-formats-title/raw-table/flextable-summary;
    # closing-title/closing-note/logo), plus background, status
    # (watermark), branding, and page-number furniture (spec section
    # 7.8) on each of them -- design.yaml sets a `background_image` and
    # presentation.yaml's `status_indicator: watermark` (spec section
    # 7.8) turns the status marker on for this build, with its text
    # falling back to `metadata.status: demo`.
    assert len(manifest["elements"]) == 16 + 4 * 6

    figure_entry = next(
        e for e in manifest["elements"] if e["slide_id"] == "concentration-time" and e["element_id"] == "figure"
    )
    assert figure_entry["type"] == "reportifyr"
    assert figure_entry["editability"] == "rendered_graphic"
    resolved_path = Path(figure_entry["resolved_path"])
    assert resolved_path.is_file()
    assert figure_entry["sha256"] == _sha256(resolved_path)
    # The whole point of this fixture: `{rpfy}:conc-time.png` resolves to
    # the actual reportifyr-produced PNG shipped in
    # examples/demo-deck/OUTPUTS/figures/, not a placeholder generated
    # for the test.
    assert figure_entry["sha256"] == _sha256(demo_deck_dir / "OUTPUTS" / "figures" / "conc-time.png")

    table_entry = next(
        e for e in manifest["elements"] if e["slide_id"] == "pk-summary" and e["element_id"] == "pk-table"
    )
    assert table_entry["type"] == "table"
    assert table_entry["editability"] == "fully_editable"
    resolved_table_path = Path(table_entry["resolved_path"])
    assert resolved_table_path.is_file()
    assert table_entry["sha256"] == _sha256(resolved_table_path)
    assert table_entry["sha256"] == _sha256(demo_deck_dir / "OUTPUTS" / "tables" / "pk-summary.csv")


# ---------------------------------------------------------------------------
# table-formats: a {rpfy}:-sourced native table next to a {rpfy}:-sourced
# .rds flextable rendered to a picture (issue #57)
# ---------------------------------------------------------------------------


def test_demo_deck_table_formats_slide_has_a_native_table_and_a_rendered_picture(
    demo_deck_dir, tmp_path, capsys
):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    table_formats_slide = list(prs.slides)[4]

    raw_table = next(
        shape for shape in table_formats_slide.shapes if shape.name == "raw-table"
    )
    assert [cell.text for cell in raw_table.table.rows[0].cells] == [
        "Participant", "Weight (kg)", "Dose (mg/kg)", "Cmax (mg/L)", "Tmax (hr)",
    ]

    flextable_picture = next(
        shape for shape in table_formats_slide.shapes if shape.name == "flextable-summary"
    )
    assert flextable_picture.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE


def test_demo_deck_table_formats_slide_footers_reflect_their_own_meta_types(
    demo_deck_dir, tmp_path, capsys
):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    table_formats_slide = list(prs.slides)[4]

    raw_table_footer = next(
        shape for shape in table_formats_slide.shapes if shape.name == "raw-table__footer"
    )
    assert "Raw per-participant values" in raw_table_footer.text_frame.text

    flextable_footer = next(
        shape for shape in table_formats_slide.shapes if shape.name == "flextable-summary__footer"
    )
    flextable_footer_text = flextable_footer.text_frame.text
    assert "Population-level summary statistics" in flextable_footer_text
    assert "Abbreviations: PK: pharmacokinetic." in flextable_footer_text


def test_demo_deck_manifest_records_the_flextable_source_not_the_rendered_png(
    demo_deck_dir, tmp_path, capsys
):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    manifest = json.loads(Path(output["manifest"]).read_text())

    flextable_entry = next(
        e for e in manifest["elements"]
        if e["slide_id"] == "table-formats" and e["element_id"] == "flextable-summary"
    )
    assert flextable_entry["type"] == "reportifyr"
    assert flextable_entry["editability"] == "rendered_graphic"
    # The manifest records the original .rds source, not the rendered PNG.
    resolved_path = Path(flextable_entry["resolved_path"])
    assert resolved_path.suffix == ".rds"
    assert resolved_path.is_file()
    assert flextable_entry["sha256"] == _sha256(resolved_path)
    assert flextable_entry["sha256"] == _sha256(
        demo_deck_dir / "OUTPUTS" / "tables" / "pk-flextable-summary.rds"
    )
    assert any(
        "rendered flextable to PNG" in warning for warning in manifest.get("warnings", [])
    )


# ---------------------------------------------------------------------------
# pk-interpretation: the two `type: quarto` elements (spec section 8.1)
# ---------------------------------------------------------------------------


def test_demo_deck_equation_fragment_renders_as_a_picture(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    pk_interpretation_slide = list(prs.slides)[2]

    figure = next(shape for shape in pk_interpretation_slide.shapes if shape.name == "figure")
    assert figure.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
    assert "elimination" in _alt_text(figure).lower()


def test_demo_deck_narrative_fragment_executes_real_r_code(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    prs = Presentation(output["output"])
    pk_interpretation_slide = list(prs.slides)[2]

    note = next(shape for shape in pk_interpretation_slide.shapes if shape.name == "note")
    text = note.text_frame.text
    # The R chunk in fragments/half-life-narrative.qmd fits a log-linear
    # elimination model per Theoph participant and reports the mean -- this
    # is a real, non-placeholder computed value, not text baked into the
    # YAML: a change to the underlying computation (or to R/knitr's own
    # numeric behavior) would change this number too.
    assert "hours" in text
    assert "population PK model" in text

    bold_runs = [
        run.text
        for paragraph in note.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.bold
    ]
    assert any("hours" in run for run in bold_runs)


def test_demo_deck_manifest_records_resolved_quarto_render_modes(demo_deck_dir, tmp_path, capsys):
    output = _build_demo_deck(demo_deck_dir, tmp_path, capsys)
    manifest = json.loads(Path(output["manifest"]).read_text())

    figure_entry = next(
        e for e in manifest["elements"]
        if e["slide_id"] == "pk-interpretation" and e["element_id"] == "figure"
    )
    assert figure_entry["type"] == "quarto"
    assert figure_entry["render_mode"] == "png"
    assert figure_entry["editability"] == "rendered_graphic"
    assert Path(figure_entry["resolved_path"]).name == "elimination-equation.qmd"

    note_entry = next(
        e for e in manifest["elements"]
        if e["slide_id"] == "pk-interpretation" and e["element_id"] == "note"
    )
    assert note_entry["type"] == "quarto"
    assert note_entry["render_mode"] == "native"
    assert note_entry["editability"] == "fully_editable"
    assert Path(note_entry["resolved_path"]).name == "half-life-narrative.qmd"
