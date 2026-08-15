import json

import pytest
import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from deckifyr.plan import expand_presentation
from deckifyr.pptx.compose import _compute_image_placement, compose_and_write
from deckifyr.schema.design import (
    BrandingFurniture,
    Defaults,
    DesignDocument,
    Fonts,
    Furniture,
    PageNumberFurniture,
    SlideSize,
    StatusFurniture,
    TableStyle,
    TextStyle,
)
from deckifyr.schema.errors import ContentValidationError
from deckifyr.schema.layouts import Box, Element, Layout, LayoutsDocument
from deckifyr.schema.presentation import (
    BuildConfig,
    DesignRef,
    Metadata,
    PresentationDocument,
    ReportifyrConfig,
    Slide,
)


def test_contain_letterboxes_and_centers():
    # 100x100 box, 200x100 (2:1) image -> width-limited, centered vertically.
    assert _compute_image_placement("contain", 100, 100, 200, 100) == (
        0, 25, 100, 50, 0.0, 0.0, 0.0, 0.0
    )


def test_cover_crops_the_long_axis():
    left, top, width, height, crop_left, crop_right, crop_top, crop_bottom = (
        _compute_image_placement("cover", 100, 100, 200, 100)
    )
    assert (left, top, width, height) == (0, 0, 100, 100)
    assert crop_left == crop_right == 0.25
    assert crop_top == crop_bottom == 0.0


def test_stretch_ignores_aspect_ratio():
    assert _compute_image_placement("stretch", 100, 50, 200, 100) == (
        0, 0, 100, 50, 0.0, 0.0, 0.0, 0.0
    )


def _design() -> DesignDocument:
    return DesignDocument(
        deckifyr="0.1",
        slide=SlideSize(width="10in", height="7.5in"),
        fonts=Fonts(body="Arial", heading="Arial"),
    )


def _presentation(*, alt_text: str | None = "a red rectangle") -> PresentationDocument:
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[
                    Element(
                        id="logo",
                        type="image",
                        source="logo.png",
                        alt_text=alt_text,
                        fit="contain",
                        box=Box(x="0in", y="0in", width="2in", height="1in"),
                    )
                ],
            )
        ],
    )


@pytest.fixture
def project(tmp_path):
    Image.new("RGB", (400, 200), color="red").save(tmp_path / "logo.png")
    (tmp_path / "design.yaml").write_text("design")
    (tmp_path / "layouts.yaml").write_text("layouts")
    (tmp_path / "presentation.yaml").write_text("presentation")
    return tmp_path


def _build(project, presentation, design):
    layouts = LayoutsDocument(deckifyr="0.1", layouts={})
    resolved = expand_presentation(presentation, design, layouts, strict=True)
    return compose_and_write(
        presentation,
        design,
        resolved,
        project_root=project,
        presentation_path=project / "presentation.yaml",
        design_path=project / "design.yaml",
        layouts_path=project / "layouts.yaml",
    )


def test_image_element_builds_and_manifest_records_its_source(project):
    result = _build(project, _presentation(), _design())

    assert result.output_path.is_file()
    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (picture,) = list(slide.shapes)
    assert picture.name == "logo"

    manifest = json.loads(result.manifest_path.read_text())
    (element_entry,) = manifest["elements"]
    assert element_entry["element_id"] == "logo"
    assert element_entry["editability"] == "rendered_graphic"
    assert element_entry["resolved_path"].endswith("logo.png")
    assert "sha256" in element_entry


def test_missing_alt_text_raises(project):
    with pytest.raises(ContentValidationError):
        _build(project, _presentation(alt_text=None), _design())


def test_slide_notes_are_written_to_the_notes_slide(project):
    presentation = _presentation()
    presentation.slides[0].notes = "Mention the Q3 numbers."
    result = _build(project, presentation, _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    assert slide.notes_slide.notes_text_frame.text == "Mention the Q3 numbers."


def test_slide_without_notes_has_no_notes_slide(project):
    result = _build(project, _presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    assert slide.has_notes_slide is False


def _shape_group_presentation() -> PresentationDocument:
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[
                    Element(
                        id="card",
                        type="group",
                        box=Box(x="0in", y="0in", width="3in", height="2in"),
                        elements=[
                            Element(
                                id="backdrop",
                                type="shape",
                                shape_kind="rounded_rectangle",
                                box=Box(x="0in", y="0in", width="3in", height="2in"),
                            ),
                            Element(
                                id="label",
                                type="text",
                                value="hello",
                                box=Box(x="0.2in", y="0.2in", width="2.6in", height="0.5in"),
                            ),
                        ],
                    )
                ],
            )
        ],
    )


def test_group_element_builds_a_native_group_with_named_children(project):
    result = _build(project, _shape_group_presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (group,) = list(slide.shapes)
    assert group.shape_type is not None
    assert group.name == "card"
    child_names = [shape.name for shape in group.shapes]
    assert child_names == ["backdrop", "label"]

    manifest = json.loads(result.manifest_path.read_text())
    element_ids = [entry["element_id"] for entry in manifest["elements"]]
    assert element_ids == ["backdrop", "label", "card"]
    editabilities = {entry["element_id"]: entry["editability"] for entry in manifest["elements"]}
    assert editabilities == {
        "backdrop": "fully_editable",
        "label": "fully_editable",
        "card": "fully_editable",
    }


def test_furniture_composes_as_ordinary_shapes(project):
    Image.new("RGB", (400, 200), color="blue").save(project / "bg.png")

    design = DesignDocument(
        deckifyr="0.1",
        slide=SlideSize(width="10in", height="7.5in", background_image="bg.png"),
        fonts=Fonts(body="Arial", heading="Arial"),
        furniture=Furniture(
            status=StatusFurniture(
                enabled=True, box=Box(x="0in", y="0in", width="1in", height="0.3in")
            ),
            branding=BrandingFurniture(
                text="Acme / R&D", box=Box(x="0in", y="0in", width="2in", height="0.3in")
            ),
            page_number=PageNumberFurniture(
                box=Box(x="0in", y="0in", width="1in", height="0.3in")
            ),
        ),
    )
    presentation = PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[Element(id="title", type="text", value="hi", box=Box(x="0in", y="0in", width="1in", height="1in"))],
            )
        ],
    )
    result = _build(project, presentation, design)

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    shape_names = {shape.name for shape in slide.shapes}
    assert shape_names == {
        "__furniture_background",
        "__furniture_status",
        "__furniture_branding",
        "__furniture_page_number",
        "title",
    }

    background = next(shape for shape in slide.shapes if shape.name == "__furniture_background")
    cnv_pr = background._element.find(qn("p:nvPicPr")).find(qn("p:cNvPr"))
    assert cnv_pr.get("descr") == "Background image"

    page_number = next(
        shape for shape in slide.shapes if shape.name == "__furniture_page_number"
    )
    assert page_number.text_frame.text == "1 / 1"


def test_shape_without_style_gets_default_outline(project):
    presentation = PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[
                    Element(
                        id="divider",
                        type="shape",
                        shape_kind="rectangle",
                        box=Box(x="0in", y="0in", width="3in", height="0.05in"),
                    )
                ],
            )
        ],
    )
    result = _build(project, presentation, _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (shape,) = list(slide.shapes)
    assert shape.name == "divider"
    assert shape.line.color.rgb == RGBColor.from_string("000000")


def _table_presentation(
    *, source: str = "data.csv", table_style: str | None = None
) -> PresentationDocument:
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(output="build/out.pptx", manifest="build/out.manifest.json"),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[
                    Element(
                        id="tbl",
                        type="table",
                        source=source,
                        table_style=table_style,
                        box=Box(x="0in", y="0in", width="4in", height="2in"),
                    )
                ],
            )
        ],
    )


def test_table_element_builds_a_native_table_from_csv(project):
    (project / "data.csv").write_text("name,score\nAda,10\nGrace,9\n")

    result = _build(project, _table_presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (graphic_frame,) = list(slide.shapes)
    assert graphic_frame.name == "tbl"
    table = graphic_frame.table
    assert [cell.text for cell in table.rows[0].cells] == ["name", "score"]
    assert [cell.text for cell in table.rows[1].cells] == ["Ada", "10"]
    assert [cell.text for cell in table.rows[2].cells] == ["Grace", "9"]

    manifest = json.loads(result.manifest_path.read_text())
    (element_entry,) = manifest["elements"]
    assert element_entry["element_id"] == "tbl"
    assert element_entry["editability"] == "fully_editable"
    assert element_entry["resolved_path"].endswith("data.csv")
    assert "sha256" in element_entry


def test_table_style_applies_fills_header_text_color_and_borders(project):
    (project / "data.csv").write_text("name,score\nAda,10\nGrace,9\n")
    design = DesignDocument(
        deckifyr="0.1",
        slide=SlideSize(width="10in", height="7.5in"),
        fonts=Fonts(body="Arial", heading="Arial"),
        colors={"primary": "#2457A6", "accent": "#D14D32"},
        table_styles={
            "branded": TableStyle(
                header_fill="primary",
                header_text_color="#FFFFFF",
                body_fill="#FFFFFF",
                band_fill="#EEF2FA",
                border_color="accent",
                border_width="1.5pt",
            )
        },
    )

    result = _build(project, _table_presentation(table_style="branded"), design)

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (graphic_frame,) = list(slide.shapes)
    table = graphic_frame.table

    header_cell = table.cell(0, 0)
    assert header_cell.fill.fore_color.rgb == RGBColor.from_string("2457A6")
    header_run = header_cell.text_frame.paragraphs[0].runs[0]
    assert header_run.font.color.rgb == RGBColor.from_string("FFFFFF")

    body_cell = table.cell(1, 0)
    assert body_cell.fill.fore_color.rgb == RGBColor.from_string("FFFFFF")
    band_cell = table.cell(2, 0)
    assert band_cell.fill.fore_color.rgb == RGBColor.from_string("EEF2FA")

    tcPr = header_cell._tc.get_or_add_tcPr()
    border_children = [qn("a:lnL"), qn("a:lnR"), qn("a:lnT"), qn("a:lnB")]
    for tag in border_children:
        line = tcPr.find(tag)
        assert line is not None
        assert line.get("w") == str(int(1.5 * 12700))
        fill_color = line.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        assert fill_color.get("val") == "D14D32"

    # Border elements must precede the fill element (OOXML schema order),
    # otherwise PowerPoint would treat the file as corrupted.
    child_tags = [child.tag for child in tcPr]
    assert child_tags.index(qn("a:lnB")) < child_tags.index(qn("a:solidFill"))


def test_table_without_table_style_keeps_default_look(project):
    (project / "data.csv").write_text("name,score\nAda,10\n")

    result = _build(project, _table_presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (graphic_frame,) = list(slide.shapes)
    table = graphic_frame.table
    header_cell = table.cell(0, 0)
    # No `table_style` set anywhere -- no explicit per-cell fill/border was
    # ever applied, so `a:tcPr` on the header cell has no fill/line children.
    tcPr = header_cell._tc.get_or_add_tcPr()
    assert len(tcPr) == 0


def test_table_source_not_found_raises(project):
    with pytest.raises(ContentValidationError):
        _build(project, _table_presentation(source="missing.csv"), _design())


def test_table_source_outside_project_root_raises(project):
    with pytest.raises(ContentValidationError):
        _build(project, _table_presentation(source="../outside/data.csv"), _design())


# ---------------------------------------------------------------------------
# Reportifyr figures + footers
# ---------------------------------------------------------------------------


def _write_reportifyr_figure_fixture(project):
    figures = project / "OUTPUTS" / "figures"
    figures.mkdir(parents=True)
    Image.new("RGB", (400, 200), color="blue").save(figures / "conc-time.png")
    metadata = {
        "source_meta": {"path": "scripts/01_analysis.R", "latest_time": "2026-08-11 22:28:55"},
        "object_meta": {
            "meta_type": "conc-time-trajectories",
            "footnotes": {
                "notes": ["Data are from the built-in Theoph dataset."],
                "abbreviations": ["PK"],
            },
        },
    }
    (figures / "conc-time_png_metadata.json").write_text(json.dumps(metadata))

    standard_footnotes = {
        "figure_footnotes": {
            "conc-time-trajectories": "This plot shows individual concentration-time trajectories."
        },
        "table_footnotes": {},
        "abbreviations": {"PK": "pharmacokinetic"},
    }
    (project / "standard_footnotes.yaml").write_text(yaml.safe_dump(standard_footnotes))


def _reportifyr_presentation(*, footer_placement: str | None = None) -> PresentationDocument:
    element_kwargs = dict(
        id="fig",
        type="reportifyr",
        value="{rpfy}:conc-time.png",
        alt_text="a concentration-time plot",
        box=Box(x="0in", y="0in", width="4in", height="2in"),
    )
    if footer_placement is not None:
        element_kwargs["footer_placement"] = footer_placement
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(
            output="build/out.pptx",
            manifest="build/out.manifest.json",
            reportifyr=ReportifyrConfig(standard_footnotes="standard_footnotes.yaml"),
        ),
        slides=[Slide(id="s1", layout=None, elements=[Element(**element_kwargs)])],
    )


def test_reportifyr_element_resolves_and_places_a_picture(project):
    _write_reportifyr_figure_fixture(project)
    result = _build(project, _reportifyr_presentation(footer_placement="none"), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (picture,) = list(slide.shapes)
    assert picture.name == "fig"

    manifest = json.loads(result.manifest_path.read_text())
    (element_entry,) = manifest["elements"]
    assert element_entry["type"] == "reportifyr"
    assert element_entry["editability"] == "rendered_graphic"
    assert element_entry["resolved_path"].endswith("conc-time.png")


def test_reportifyr_footer_below_places_a_shape_beneath_the_box(project):
    _write_reportifyr_figure_fixture(project)
    result = _build(project, _reportifyr_presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    picture, footer = list(slide.shapes)
    assert footer.name == "fig__footer"
    assert footer.top == picture.top + picture.height

    footer_text = footer.text_frame.text
    assert footer_text.startswith("Source: scripts/01_analysis.R")
    assert "Notes: This plot shows individual concentration-time trajectories." in footer_text
    assert "Abbreviations: PK: pharmacokinetic." in footer_text
    assert slide.has_notes_slide is False


def test_reportifyr_footer_inherits_every_field_of_a_named_style(project):
    _write_reportifyr_figure_fixture(project)
    design = DesignDocument(
        deckifyr="0.1",
        slide=SlideSize(width="10in", height="7.5in"),
        fonts=Fonts(body="Arial", heading="Arial"),
        text_styles={
            "loud-footnote": TextStyle(
                font="body", size="9pt", bold=True, italic=True, color="#FF0000"
            )
        },
        defaults=Defaults(footer_style="loud-footnote"),
    )
    result = _build(project, _reportifyr_presentation(), design)

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    _picture, footer = list(slide.shapes)
    run = footer.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.italic is True
    assert run.font.size.pt == 9
    assert run.font.color.rgb == RGBColor.from_string("FF0000")


def test_reportifyr_footer_notes_appends_to_slide_notes(project):
    _write_reportifyr_figure_fixture(project)
    result = _build(project, _reportifyr_presentation(footer_placement="notes"), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    (picture,) = list(slide.shapes)
    assert picture.name == "fig"
    assert "Abbreviations: PK: pharmacokinetic." in slide.notes_slide.notes_text_frame.text


def test_reportifyr_footer_none_skips_footer_entirely(project):
    _write_reportifyr_figure_fixture(project)
    result = _build(project, _reportifyr_presentation(footer_placement="none"), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    assert len(list(slide.shapes)) == 1
    assert slide.has_notes_slide is False


def test_reportifyr_missing_sidecar_raises(project):
    figures = project / "OUTPUTS" / "figures"
    figures.mkdir(parents=True)
    Image.new("RGB", (400, 200), color="blue").save(figures / "conc-time.png")
    (project / "standard_footnotes.yaml").write_text(yaml.safe_dump({"abbreviations": {}}))

    with pytest.raises(ContentValidationError):
        _build(project, _reportifyr_presentation(), _design())


def test_reportifyr_element_requires_alt_text(project):
    _write_reportifyr_figure_fixture(project)
    presentation = _reportifyr_presentation(footer_placement="none")
    presentation.slides[0].elements[0].alt_text = None
    with pytest.raises(ContentValidationError):
        _build(project, presentation, _design())


def _rpfy_table_presentation() -> PresentationDocument:
    return PresentationDocument(
        deckifyr="0.1",
        design=DesignRef(base="design.yaml"),
        layouts="layouts.yaml",
        metadata=Metadata(title="Test"),
        build=BuildConfig(
            output="build/out.pptx",
            manifest="build/out.manifest.json",
            reportifyr=ReportifyrConfig(standard_footnotes="standard_footnotes.yaml"),
        ),
        slides=[
            Slide(
                id="s1",
                layout=None,
                elements=[
                    Element(
                        id="tbl",
                        type="table",
                        source="{rpfy}:pk-summary.csv",
                        box=Box(x="0in", y="0in", width="4in", height="2in"),
                    )
                ],
            )
        ],
    )


def test_rpfy_sourced_table_builds_and_footers_from_table_footnotes(project):
    tables = project / "OUTPUTS" / "tables"
    tables.mkdir(parents=True)
    (tables / "pk-summary.csv").write_text("name,score\nAda,10\n")
    metadata = {
        "object_meta": {
            "meta_type": "univariate",
            "footnotes": {"notes": [], "abbreviations": []},
        }
    }
    (tables / "pk-summary_csv_metadata.json").write_text(json.dumps(metadata))
    standard_footnotes = {
        "table_footnotes": {"univariate": "The p-value is from the likelihood ratio test."},
        "figure_footnotes": {},
        "abbreviations": {},
    }
    (project / "standard_footnotes.yaml").write_text(yaml.safe_dump(standard_footnotes))

    result = _build(project, _rpfy_table_presentation(), _design())

    prs = Presentation(str(result.output_path))
    (slide,) = list(prs.slides)
    graphic_frame, footer = list(slide.shapes)
    assert graphic_frame.table.rows[1].cells[0].text == "Ada"
    assert footer.text_frame.text == "Notes: The p-value is from the likelihood ratio test."

    manifest = json.loads(result.manifest_path.read_text())
    (element_entry,) = manifest["elements"]
    assert element_entry["resolved_path"].endswith("pk-summary.csv")
